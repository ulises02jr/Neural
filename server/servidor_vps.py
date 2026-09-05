"""
servidor_vps.py - Servidor Flask para modo ensayo (VPS público).

Diferente del sistema live:
- No usa MIDI ni Logic
- Cualquier músico puede ver cualquier canción cuando quiera
- Login con 2 passwords: uno para músicos (lectura), uno para admin (subir JSONs)
- Página principal: setlist destacado + biblioteca completa
- Página admin: subir JSONs, armar setlist, eliminar canciones

Endpoints:
  /                       Setlist + biblioteca (requiere login músicos)
  /cancion/<numero>       Ver una canción específica
  /login                  Login músicos
  /admin                  Panel admin (requiere login admin)
  /admin/login            Login admin
  /admin/subir            POST: subir un JSON nuevo
  /admin/eliminar/<n>     POST: eliminar canción
  /admin/setlist          POST: actualizar el setlist
  /api/cancion/<n>        JSON crudo de una canción (requiere login)
  /logout                 Cerrar sesión

Configuración: archivo config.json con:
  - password_musicos (sha256 hash)
  - password_admin (sha256 hash)
  - secret_key (cookie de sesión)
  - setlist (lista de números)
"""

import os
import json
import re
import subprocess
import shutil
import threading
import logging
import hashlib
import hmac
import secrets
from functools import wraps
from pathlib import Path
from datetime import timedelta

from flask import (
    Flask, render_template, request, redirect, url_for,
    session, flash, jsonify, abort, send_file
)
from werkzeug.utils import secure_filename

from transposicion import transponer_cancion, NOTA_A_INDICE
import usuarios
import emails as emails_module


# ───────────────────────── Configuración base ─────────────────────────
BASE_DIR = Path(__file__).resolve().parent
CARPETA_ORGS = BASE_DIR / "orgs"
CARPETA_CANCIONES = BASE_DIR / "canciones"
ARCHIVO_CONFIG = BASE_DIR / "config.json"

CARPETA_CANCIONES.mkdir(exist_ok=True)


# ───────────────────── Almacenamiento por organización ─────────────────────
# Org #1 (instalación original) usa las carpetas actuales tal cual (no se mueven).
# Organizaciones nuevas viven en BASE_DIR/orgs/<id>/{canciones,pistas,pads,config.json}.
def _org_base(org=1):
    try:
        if org and int(org) != 1:
            return CARPETA_ORGS / str(int(org))
    except Exception:
        pass
    return BASE_DIR


def dir_canciones(org=1):
    d = _org_base(org) / "canciones"
    d.mkdir(parents=True, exist_ok=True)
    return d


def dir_pistas(org=1):
    d = _org_base(org) / "pistas"
    d.mkdir(parents=True, exist_ok=True)
    return d


def dir_pads(org=1):
    d = _org_base(org) / "pads"
    d.mkdir(parents=True, exist_ok=True)
    return d


def archivo_config(org=1):
    try:
        if org and int(org) != 1:
            base = CARPETA_ORGS / str(int(org))
            base.mkdir(parents=True, exist_ok=True)
            return base / "config.json"
    except Exception:
        pass
    return ARCHIVO_CONFIG


def hash_password(plain):
    return hashlib.sha256(plain.encode("utf-8")).hexdigest()


def cargar_config(org=1):
    """Carga el config.json de la organización. Si no existe, lo crea con defaults."""
    _cfg_path = archivo_config(org)
    if not _cfg_path.exists():
        # Primera ejecución: crear config con passwords por defecto
        default = {
            "password_musicos": hash_password("musicos2026"),
            "password_admin": hash_password("admin2026"),
            "secret_key": secrets.token_hex(32),
            "live_token": secrets.token_hex(16),
            "setlists": [],  # Multi-setlists: lista de {id, nombre, fecha, canciones}
            "live_activo": False,
            "mac_local_ip": None,
            "ultimo_heartbeat": None,
        }
        with open(_cfg_path, "w", encoding="utf-8") as f:
            json.dump(default, f, indent=2)
        print("⚠️  config.json creado con passwords por defecto.")
        print("    Cámbialos editando el archivo.")
        return default
    with open(_cfg_path, "r", encoding="utf-8") as f:
        cfg = json.load(f)
    # Migración automática: agregar campos nuevos si la config es vieja
    cambios = False
    if "live_token" not in cfg:
        cfg["live_token"] = secrets.token_hex(16)
        cambios = True
    if "live_activo" not in cfg:
        cfg["live_activo"] = False
        cambios = True
    if "mac_local_ip" not in cfg:
        cfg["mac_local_ip"] = None
        cambios = True
    if "ultimo_heartbeat" not in cfg:
        cfg["ultimo_heartbeat"] = None
        cambios = True
    if "org_nombre" not in cfg:
        cfg["org_nombre"] = "Mi Iglesia Internacional"
        cambios = True
    # Migración: setlist único → multi-setlists
    if "setlists" not in cfg:
        cfg["setlists"] = []
        # Si había un setlist viejo, convertirlo
        viejo = cfg.pop("setlist", None)
        if viejo:
            cfg["setlists"].append({
                "id": secrets.token_hex(6),
                "nombre": "Setlist migrado",
                "fecha": _hoy_iso(),
                "canciones": [
                    {"id": n, "tono": None} if isinstance(n, int)
                    else {"id": int(n.get("id", 0)), "tono": n.get("tono")}
                    for n in viejo if (isinstance(n, int) or isinstance(n, dict))
                ],
                "creado": _ahora_iso(),
            })
            print(f"✓ Setlist viejo migrado a multi-setlists")
        cambios = True
    # (Los repertorios ya NO se borran automáticamente por fecha; el usuario los elimina cuando quiera.)
    if cambios:
        with open(_cfg_path, "w", encoding="utf-8") as f:
            json.dump(cfg, f, indent=2)
    return cfg


def _hoy_iso():
    """Fecha de hoy en formato YYYY-MM-DD."""
    from datetime import date
    return date.today().isoformat()


def _ahora_iso():
    """Timestamp ISO."""
    from datetime import datetime
    return datetime.now().isoformat(timespec="seconds")


def guardar_config(cfg, org=None):
    org = org if org is not None else org_actual()
    with open(archivo_config(org), "w", encoding="utf-8") as f:
        json.dump(cfg, f, indent=2)


def _crear_config_org(org_id, token, org_nombre):
    """Crea el config.json inicial de una organización nueva.
    Clave: live_token = token de la organización (para que NeuralPlay/API funcione)."""
    cfg = {
        "password_musicos": secrets.token_hex(32),
        "password_admin": secrets.token_hex(32),
        "secret_key": secrets.token_hex(32),
        "live_token": token,
        "setlists": [],
        "live_activo": False,
        "mac_local_ip": None,
        "ultimo_heartbeat": None,
        "org_nombre": org_nombre,
    }
    guardar_config(cfg, org=org_id)
    return cfg


# ───────────────────────── App Flask ─────────────────────────
# Cargamos config una vez para inicializar Flask
_config_inicial = cargar_config()

# Inicializar base de datos de usuarios (crea tablas si no existen)
usuarios.init_db()

app = Flask(__name__)
app.secret_key = _config_inicial["secret_key"]
app.permanent_session_lifetime = timedelta(days=30)
app.config["MAX_CONTENT_LENGTH"] = 200 * 1024 * 1024  # 5 MB max por archivo
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "Lax"
app.config["SESSION_COOKIE_SECURE"] = True


def _client_ip():
    xff = request.headers.get("X-Forwarded-For", "")
    if xff:
        return xff.split(",")[0].strip()
    return request.remote_addr or "?"


@app.context_processor
def _inyectar_csrf():
    if "_csrf" not in session:
        session["_csrf"] = secrets.token_hex(32)
    return {"csrf_token": lambda: session.get("_csrf", "")}


@app.before_request
def _proteger_csrf():
    if request.method in ("GET", "HEAD", "OPTIONS"):
        return
    if request.path.startswith("/api/"):
        return
    token = session.get("_csrf")
    enviado = request.form.get("_csrf") or request.headers.get("X-CSRFToken")
    if not token or not enviado or not hmac.compare_digest(str(token), str(enviado)):
        abort(400)


@app.context_processor
def inject_org():
    try:
        nombre = get_config().get("org_nombre", "Mi Iglesia Internacional")
    except Exception:
        nombre = "Mi Iglesia Internacional"
    try:
        logo_ver = int((BASE_DIR / "static" / "logo.png").stat().st_mtime)
    except Exception:
        logo_ver = 0
    return {"org_nombre": nombre, "logo_ver": logo_ver}


@app.context_processor
def inject_acento():
    """#2 Color de acento por usuario (global por cuenta)."""
    try:
        uid = session.get("user_id")
        ac = usuarios.obtener_acento(uid) if uid else None
    except Exception:
        ac = None
    return {"mi_acento": ac or ""}


@app.route("/api/perfil/acento", methods=["GET", "POST"])
def api_perfil_acento():
    """#2 Lee o guarda el color de acento del usuario logueado."""
    uid = session.get("user_id")
    if not uid:
        return jsonify({"ok": False, "error": "no_session"}), 401
    if request.method == "POST":
        if request.is_json:
            color = (request.get_json(silent=True) or {}).get("acento")
        else:
            color = request.form.get("acento")
        ok = usuarios.guardar_acento(uid, color)
        return jsonify({"ok": bool(ok), "acento": usuarios.obtener_acento(uid) or ""})
    return jsonify({"ok": True, "acento": usuarios.obtener_acento(uid) or ""})


@app.context_processor
def inject_prefs():
    """Preferencias de vista del visor por usuario (global por cuenta)."""
    try:
        uid = session.get("user_id")
        pr = usuarios.obtener_prefs(uid) if uid else "{}"
    except Exception:
        pr = "{}"
    return {"mi_prefs": pr or "{}"}


@app.route("/api/perfil/prefs", methods=["GET", "POST"])
def api_perfil_prefs():
    """Lee o guarda las preferencias de vista del visor del usuario logueado."""
    uid = session.get("user_id")
    if not uid:
        return jsonify({"ok": False, "error": "no_session"}), 401
    if request.method == "POST":
        data = request.get_json(silent=True) if request.is_json else request.form.to_dict()
        ok = usuarios.guardar_prefs(uid, data or {})
        return jsonify({"ok": bool(ok)})
    return jsonify({"ok": True, "prefs": usuarios.obtener_prefs(uid)})


def get_config(org=None):
    """Lee el config.json (de la organización actual) del disco en cada llamada.
    Esto garantiza que múltiples workers vean los cambios."""
    org = org if org is not None else org_actual()
    return cargar_config(org)


# ───────────────────────── Helpers ─────────────────────────
def cargar_biblioteca(org=None):
    """Lee todos los JSONs de canciones/ (de la organización) → dict numero → datos."""
    org = _org_req(org)
    biblioteca = {}
    for archivo in dir_canciones(org).glob("*.json"):
        try:
            with open(archivo, "r", encoding="utf-8") as f:
                datos = json.load(f)
                if "numero" in datos:
                    biblioteca[datos["numero"]] = datos
        except Exception as e:
            print(f"⚠️  Error leyendo {archivo.name}: {e}")
    return biblioteca


def login_required(rol):
    """Decorador. rol='musico' o 'admin'.
    Compatible con:
    - Sistema viejo: session['rol'] = 'musico' o 'admin' (password compartido)
    - Sistema nuevo: session['user_id'] + session['rol'] + session['nombre']
    Para 'musico', admin también puede acceder.
    """
    def wrapper(fn):
        @wraps(fn)
        def decorated(*args, **kwargs):
            if session.get("rol") != rol:
                # Admin puede acceder a rutas de musico
                if rol == "musico" and session.get("rol") == "admin":
                    return fn(*args, **kwargs)
                return redirect(url_for("login" if rol == "musico" else "admin_login"))
            return fn(*args, **kwargs)
        return decorated
    return wrapper


def get_usuario_actual():
    """Devuelve dict del usuario logueado o None.
    Si es login viejo (password compartido), devuelve un dict mínimo.
    """
    uid = session.get("user_id")
    if uid:
        u = usuarios.buscar_por_id(uid)
        if u:
            return u
    # Fallback: login con password compartido (admin de emergencia)
    if session.get("rol") == "admin" and session.get("nombre") == "Admin (Emergencia)":
        return {"nombre": "Admin", "apellido": "(Emergencia)", "email": "—", "rol": "admin", "id": None}
    return None


def org_actual():
    """Organización de la request actual.
    - Web: por la sesión (org_id guardado al iniciar sesión).
    - API/NeuralPlay: por el token (live_token → organización dueña).
    - Fallback: organización #1 (legado / instalación original).
    NOTA: solo llamar dentro de un contexto de request; los hilos de render
    reciben el org_id explícitamente, no usan este helper.
    """
    try:
        oid = session.get("org_id")
        if oid:
            return oid
        auth = request.headers.get("Authorization", "")
        token = auth[7:].strip() if auth[:7].lower() == "bearer " else request.args.get("token", "")
        if token:
            o = usuarios.obtener_org_por_token(token)
            if o:
                return o["id"]
    except Exception:
        pass
    return 1


_org_ctx = threading.local()


def _set_org_ctx(org):
    """Fija la organización 'ambiente' del hilo actual (para hilos de render)."""
    try:
        _org_ctx.value = int(org)
    except Exception:
        _org_ctx.value = 1


def _cur_org():
    """Organización ambiente: thread-local (hilo de render) o la de la request."""
    v = getattr(_org_ctx, "value", None)
    return v if v else org_actual()


def _org_req(org=None):
    """Resuelve la organización: si se pasa explícita, la usa; si no, la ambiente."""
    return org if org is not None else _cur_org()


# ───────────────────────── Rutas públicas (login) ─────────────────────────
@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")
        if not email or not password:
            flash("Completá email y contraseña", "error")
            return render_template("login_musicos.html")
        ip = _client_ip()
        if usuarios.login_bloqueado(ip):
            flash("Demasiados intentos fallidos. Esperá unos minutos y volvé a intentar.", "error")
            return render_template("login_musicos.html")
        # Intentar autenticar con sistema de usuarios
        u = usuarios.autenticar(email, password)
        if u:
            # Si es admin, redirigir al login de admin (separación de pantallas)
            if u["rol"] == "admin":
                flash("Sos administrador. Usá el acceso de administrador.", "error")
                return redirect(url_for("admin_login"))
            usuarios.limpiar_intentos(ip)
            session.permanent = True
            session.clear()
            session["user_id"] = u["id"]
            session["rol"] = u["rol"]
            session["nombre"] = u["nombre"]
            session["org_id"] = u.get("org_id") or 1
            return redirect(url_for("principal"))
        # Verificar si el usuario existe pero está pendiente
        existente = usuarios.buscar_por_email(email)
        if existente and existente["estado"] == "pendiente":
            flash("Tu cuenta está pendiente de aprobación", "error")
        elif existente and existente["estado"] == "rechazado":
            flash("Tu cuenta no fue aprobada. Contactá al administrador.", "error")
        else:
            usuarios.registrar_intento(ip)
            flash("Email o contraseña incorrectos", "error")
    return render_template("login_musicos.html")


@app.route("/registro", methods=["GET", "POST"])
def registro():
    # Registro abierto cerrado: los músicos se unen a una organización con su código.
    return redirect(url_for("unirse"))


@app.route("/unirse", methods=["GET", "POST"])
def unirse():
    """Un músico se une a una organización con su código. Queda PENDIENTE de aprobación."""
    if request.method == "POST":
        codigo = request.form.get("codigo", "").strip().upper()
        nombre = request.form.get("nombre", "").strip()
        apellido = request.form.get("apellido", "").strip()
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")
        password2 = request.form.get("password2", "")
        prev = dict(codigo=codigo, nombre=nombre, apellido=apellido, email=email)
        if not codigo or not nombre or not apellido or not email or not password:
            flash("Completá todos los campos", "error")
            return render_template("unirse.html", **prev)
        if password != password2:
            flash("Las contraseñas no coinciden", "error")
            return render_template("unirse.html", **prev)
        org = usuarios.buscar_org_por_codigo(codigo)
        if not org:
            flash("Código de organización inválido", "error")
            return render_template("unirse.html", **prev)
        ok, res = usuarios.crear_usuario(nombre, apellido, email, password,
                                         rol="musico", estado="pendiente", org_id=org["id"])
        if not ok:
            flash(res, "error")
            return render_template("unirse.html", **prev)
        flash("✓ Cuenta creada en «%s». Esperá la aprobación del administrador para ingresar." % org["nombre"], "success")
        return redirect(url_for("login"))
    return render_template("unirse.html", codigo=request.args.get("codigo", "").strip().upper())


@app.route("/crear-organizacion", methods=["GET", "POST"])
def crear_organizacion():
    """Registro de una organización nueva (iglesia) + su admin dueño (auto-activo)."""
    if request.method == "POST":
        org_nombre = request.form.get("org_nombre", "").strip()
        nombre = request.form.get("nombre", "").strip()
        apellido = request.form.get("apellido", "").strip()
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")
        password2 = request.form.get("password2", "")
        prev = dict(org_nombre=org_nombre, nombre=nombre, apellido=apellido, email=email)
        if not org_nombre or not nombre or not apellido or not email or not password:
            flash("Completá todos los campos", "error")
            return render_template("crear_organizacion.html", **prev)
        if password != password2:
            flash("Las contraseñas no coinciden", "error")
            return render_template("crear_organizacion.html", **prev)
        if len(password) < 6:
            flash("La contraseña debe tener al menos 6 caracteres", "error")
            return render_template("crear_organizacion.html", **prev)
        ok, res = usuarios.crear_org_con_dueno(org_nombre, nombre, apellido, email, password)
        if not ok:
            flash(res, "error")
            return render_template("crear_organizacion.html", **prev)
        # Crear el config.json de la organización con su token
        try:
            _crear_config_org(res["org_id"], res["token"], org_nombre)
        except Exception as e:
            logging.error("crear config org %s: %s", res.get("org_id"), e)
        # Iniciar sesión como admin dueño de la nueva organización
        session.permanent = True
        session.clear()
        session["user_id"] = res["user_id"]
        session["rol"] = "admin"
        session["nombre"] = nombre
        session["org_id"] = res["org_id"]
        flash("✓ ¡Organización creada! Bienvenido a NeuralWorship.", "success")
        return redirect(url_for("admin"))
    return render_template("crear_organizacion.html")


@app.route("/admin/login", methods=["GET", "POST"])
def admin_login():
    """Login para admins. Soporta:
    - Email + password (sistema nuevo)
    - Password de emergencia (fallback para no quedar bloqueado)
    """
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        password = request.form.get("password", "")

        ip = _client_ip()
        if usuarios.login_bloqueado(ip):
            flash("Demasiados intentos fallidos. Esperá unos minutos.", "error")
            return render_template("login_admin.html")

        # Si dejan email vacío → intentar password de emergencia
        if not email and password:
            cfg = get_config()
            if hash_password(password) == cfg.get("password_admin", ""):
                usuarios.limpiar_intentos(ip)
                session.permanent = True
                session.clear()
                session["rol"] = "admin"
                session["nombre"] = "Admin (Emergencia)"
                session["org_id"] = 1
                flash("⚠️ Entraste con password de emergencia. Iniciá sesión con tu cuenta personal cuando puedas.", "success")
                return redirect(url_for("admin"))
            usuarios.registrar_intento(ip)
            flash("Password de emergencia incorrecto", "error")
            return render_template("login_admin.html")

        # Login normal con email
        u = usuarios.autenticar(email, password)
        if u and u["rol"] == "admin":
            usuarios.limpiar_intentos(ip)
            session.permanent = True
            session.clear()
            session["user_id"] = u["id"]
            session["rol"] = "admin"
            session["nombre"] = u["nombre"]
            session["org_id"] = u.get("org_id") or 1
            return redirect(url_for("admin"))
        usuarios.registrar_intento(ip)
        flash("Email o contraseña incorrectos (o no sos admin)", "error")
    return render_template("login_admin.html")


@app.route("/olvide_password", methods=["GET", "POST"])
def olvide_password():
    """Solicitar código de reset por email."""
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        if not email:
            flash("Ingresá tu email", "error")
            return render_template("olvide_password.html")
        u = usuarios.buscar_por_email(email)
        # Por seguridad: siempre damos el mismo mensaje aunque no exista
        if u and u["estado"] == "activo":
            codigo = usuarios.crear_codigo_reset(u["id"])
            nombre_completo = f"{u['nombre']} {u['apellido']}"
            ok, _ = emails_module.enviar_email_codigo_reset(u["email"], nombre_completo, codigo)
            if not ok:
                print(f"⚠️  No se pudo enviar email a {email}")
        flash("Si el email existe en nuestro sistema, te enviamos un código de reset.", "success")
        return redirect(url_for("reset_password", email=email))
    return render_template("olvide_password.html")


@app.route("/reset_password", methods=["GET", "POST"])
def reset_password():
    """Ingresar código + nueva contraseña."""
    email_pre = request.args.get("email", "")
    if request.method == "POST":
        email = request.form.get("email", "").strip().lower()
        codigo = request.form.get("codigo", "").strip()
        nueva = request.form.get("password", "")
        nueva2 = request.form.get("password2", "")
        if nueva != nueva2:
            flash("Las contraseñas no coinciden", "error")
            return render_template("reset_password.html", email=email, codigo=codigo)
        ok, msg = usuarios.usar_codigo_y_cambiar_password(email, codigo, nueva)
        if ok:
            flash("✓ Contraseña cambiada. Ya podés ingresar.", "success")
            return redirect(url_for("login"))
        flash(msg, "error")
        return render_template("reset_password.html", email=email, codigo=codigo)
    return render_template("reset_password.html", email=email_pre, codigo="")


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


# ───────────────────────── Rutas para músicos ─────────────────────────
@app.route("/")
@login_required("musico")
def principal():
    biblioteca = cargar_biblioteca()
    biblioteca_ordenada = sorted(biblioteca.values(), key=lambda c: c["numero"])
    cfg = get_config()

    # Ordenar setlists por fecha ASC (más próximo primero), luego por created DESC para mismo día
    setlists_raw = cfg.get("setlists", [])
    setlists_ordenados = sorted(
        setlists_raw,
        key=lambda s: (s.get("fecha", "9999"), s.get("creado", ""))
    )

    # Setlist seleccionado: ?setlist_id=xxx o el más próximo
    setlist_id = request.args.get("setlist_id")
    setlist_activo = None
    if setlist_id:
        setlist_activo = next((s for s in setlists_ordenados if s["id"] == setlist_id), None)
    if not setlist_activo and setlists_ordenados:
        setlist_activo = setlists_ordenados[0]  # el más próximo

    # Resolver canciones del setlist activo (con tono override)
    setlist_canciones = []
    if setlist_activo:
        for item in setlist_activo.get("canciones", []):
            num = item.get("id")
            if num in biblioteca:
                cancion = dict(biblioteca[num])
                # Aplicar tono override
                if item.get("tono"):
                    cancion["tono"] = item["tono"]
                cancion["_es_override"] = item.get("tono") and item["tono"] != biblioteca[num].get("tono")
                setlist_canciones.append(cancion)

    # Lista resumida de TODOS los setlists para el selector
    setlists_resumen = [
        {"id": s["id"], "nombre": s.get("nombre", "Sin nombre"), "fecha": s.get("fecha", "")}
        for s in setlists_ordenados
    ]

    # Estado del live
    live_activo = is_live_activo(cfg)

    return render_template(
        "principal.html",
        setlist=setlist_canciones,
        setlist_activo=setlist_activo,
        setlists_disponibles=setlists_resumen,
        biblioteca=biblioteca_ordenada,
        es_admin=(session.get("rol") == "admin"),
        live_activo=live_activo,
        live_ip=cfg.get("mac_local_ip") if live_activo else None,
        usuario=get_usuario_actual(),
    )


@app.route("/cancion/<int:numero>")
@login_required("musico")
def ver_cancion(numero):
    biblioteca = cargar_biblioteca()
    if numero not in biblioteca:
        abort(404)

    cancion = dict(biblioteca[numero])  # copia para no mutar la biblioteca
    cfg = get_config()

    # Si viene de un setlist específico, aplicar el tono override de ese setlist
    setlist_id = request.args.get("setlist_id")
    if setlist_id:
        sl = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
        if sl:
            item = next((c for c in sl.get("canciones", []) if c.get("id") == numero), None)
            if item and item.get("tono") and item["tono"] != cancion.get("tono"):
                # Calcular semitonos entre tono original y override
                from transposicion import NOTA_A_INDICE
                import re
                def raiz(t):
                    m = re.match(r'^([A-G][#b]?)', t or "")
                    return m.group(1) if m else None
                r_orig = raiz(cancion.get("tono", ""))
                r_dest = raiz(item["tono"])
                if r_orig and r_dest and r_orig in NOTA_A_INDICE and r_dest in NOTA_A_INDICE:
                    sem = (NOTA_A_INDICE[r_dest] - NOTA_A_INDICE[r_orig]) % 12
                    if sem != 0:
                        cancion = transponer_cancion(cancion, sem)

    # Soporte de transposición manual (botones +/-): ?t=N
    # Y origen: ?from=bib indica que se accedió desde biblioteca (se muestran botones +/-)
    try:
        semitonos = int(request.args.get("t", 0))
        if semitonos < -12 or semitonos > 12:
            semitonos = 0
    except (ValueError, TypeError):
        semitonos = 0

    tono_original = cancion.get("tono", "C")

    if semitonos != 0:
        cancion = transponer_cancion(cancion, semitonos)

    # Permitir transposición solo si viene de biblioteca (no desde setlist)
    desde_biblioteca = request.args.get("from") == "bib"

    if desde_biblioteca:
        volver_url = url_for("principal") + "?tab=biblioteca"
    elif setlist_id:
        volver_url = url_for("principal", setlist_id=setlist_id)
    else:
        volver_url = url_for("principal")

    return render_template(
        "visor.html",
        cancion=cancion,
        es_admin=(session.get("rol") == "admin"),
        permitir_transposicion=desde_biblioteca,
        semitonos_actual=semitonos,
        tono_original=tono_original,
        volver_url=volver_url,
    )


@app.route("/api/sync/biblioteca")
def api_sync_biblioteca():
    """Para el Puente local (app de escritorio): entrega toda la biblioteca de
    canciones (JSON completo). Autenticacion por token (server-to-server)."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    bib = cargar_biblioteca()
    canciones = [bib[n] for n in sorted(bib.keys())]
    # Repertorio activo (el mas proximo por fecha) para el Puente local
    cfg = get_config()
    setlists = sorted(cfg.get("setlists", []),
                      key=lambda s: (s.get("fecha", "9999"), s.get("creado", "")))
    setlist_out = []
    setlist_nombre = ""
    if setlists:
        activo = setlists[0]
        setlist_nombre = activo.get("nombre", "")
        for it in activo.get("canciones", []):
            nid = it.get("id")
            if nid in bib:
                setlist_out.append({"id": nid, "tono": it.get("tono")})
    return jsonify({"ok": True, "total": len(canciones), "canciones": canciones,
                    "setlist": setlist_out, "setlist_nombre": setlist_nombre})


def _token_ok():
    auth = request.headers.get("Authorization", "")
    token = auth[7:].strip() if auth[:7].lower() == "bearer " else request.args.get("token", "")
    if not token:
        return False
    try:
        if usuarios.obtener_org_por_token(token):
            return True
    except Exception:
        pass
    return token == get_config(1).get("live_token")


def _semitono_override(base_tono, override_tono):
    """Semitono firmado (mas cercano) del override respecto al tono original. 0 si igual/None."""
    import re
    from transposicion import NOTA_A_INDICE
    def raiz(t):
        m = re.match(r"^([A-G][#b]?)", t or "")
        return m.group(1) if m else None
    ro, rd = raiz(base_tono), raiz(override_tono)
    if not ro or not rd or ro not in NOTA_A_INDICE or rd not in NOTA_A_INDICE:
        return 0
    d = (NOTA_A_INDICE[rd] - NOTA_A_INDICE[ro]) % 12
    if d > 6:
        d -= 12
    return d


@app.route("/api/live/setlists")
def api_live_setlists():
    """NeuralPlay: lista de repertorios (setlists) + indice liviano de canciones. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    cfg = get_config()
    bib = cargar_biblioteca()
    idx = {}
    for n, c in bib.items():
        idx[str(n)] = {"id": n, "titulo": c.get("titulo", ""), "artista": c.get("artista", ""),
                       "tono": c.get("tono", ""), "tempo": c.get("tempo"),
                       "compas": c.get("compas", ""), "portada": c.get("portada", ""),
                       "portada_ts": c.get("portada_ts")}
    setlists_out = []
    for s in cfg.get("setlists", []):
        so = dict(s)
        cs = []
        for it in s.get("canciones", []):
            base = bib.get(it.get("id"), {})
            override = it.get("tono")
            sem = _semitono_override(base.get("tono", ""), override) if override else 0
            ci = dict(it)
            ci["tono_semitonos"] = sem
            ci["tono_nombre"] = override or base.get("tono", "")
            cs.append(ci)
        so["canciones"] = cs
        setlists_out.append(so)
    return jsonify({"ok": True, "setlists": setlists_out, "canciones": idx})


@app.route("/api/live/pistas/<int:numero>")
def api_live_pistas(numero):
    """NeuralPlay: stems + tiempos de seccion de una cancion a un tono. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    try:
        n = int(request.args.get("t", "0"))
    except ValueError:
        n = 0
    listo = _tono_listo(numero, n)
    stems = []
    fam = _leer_familias(numero)
    if listo:
        for f in sorted(_carpeta_tono(numero, n).iterdir()):
            if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO:
                stems.append({"name": f.stem, "file": f.name,
                              "familia": fam.get(f.name) or _familia_auto(f.stem)})
    cancion = cargar_biblioteca().get(numero, {})
    return jsonify({"numero": numero, "tono": n, "listo": listo,
                    "hay_pistas": len(_stems_originales(numero)) > 0,
                    "tempo": cancion.get("tempo"), "compas": cancion.get("compas", ""),
                    "stems": stems, "secciones": _leer_secciones(numero)})


@app.route("/api/live/pista/<int:numero>/<path:archivo>")
def api_live_pista(numero, archivo):
    """NeuralPlay: descarga de un stem. Token."""
    if not _token_ok():
        abort(403)
    try:
        n = int(request.args.get("t", "0"))
    except ValueError:
        n = 0
    base = _carpeta_stems(numero, n).resolve()
    ruta = (base / archivo).resolve()
    try:
        dentro = os.path.commonpath([str(base), str(ruta)]) == str(base)
    except ValueError:
        dentro = False
    if not dentro or not ruta.is_file():
        abort(404)
    return send_file(str(ruta))


@app.route("/api/live/midi/<int:numero>")
def api_live_midi(numero):
    """NeuralPlay: cajas MIDI + notas (en segundos) de una cancion. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    midi = _leer_midi(numero)
    cajas = []
    for cid, nom, chdef in MIDI_CAJAS:
        c = midi["cajas"].get(cid, {})
        notas = []
        for e in c.get("notas", []):
            if e.get("desactivar"):
                continue
            seg = _parse_tiempo(str(e.get("t", "")))
            if seg is None:
                continue
            try:
                vel = int(e.get("velocidad", 100))
            except Exception:
                vel = 100
            notas.append({
                "seg": round(seg, 3),
                "note": _midi_note_num(e.get("nota", "C"), e.get("octava", 0)),
                "vel": max(1, min(127, vel)),
                "desc": str(e.get("descripcion", "")),
            })
        notas.sort(key=lambda x: x["seg"])
        cajas.append({"id": cid, "nombre": nom, "canal": int(c.get("canal", chdef)), "notas": notas})
    return jsonify({"ok": True, "numero": numero, "cajas": cajas})


@app.route("/api/live/chart/<int:numero>")
def api_live_chart(numero):
    """NeuralPlay: chart (secciones con acordes+letra) de una cancion para servir a los musicos. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    c = cargar_biblioteca().get(numero)
    if not c:
        return jsonify({"ok": False, "error": "not found"}), 404
    try:
        sem = int(request.args.get("t", 0))
    except (ValueError, TypeError):
        sem = 0
    if sem != 0:
        import copy as _copy
        c = transponer_cancion(_copy.deepcopy(c), sem)
    return jsonify({
        "ok": True,
        "numero": numero,
        "titulo": c.get("titulo", ""),
        "artista": c.get("artista", ""),
        "tono": c.get("tono", ""),
        "tempo": c.get("tempo"),
        "compas": c.get("compas", ""),
        "secciones": c.get("secciones", []),
    })


def _beatgrid_musical(numero):
    """Grilla de negras (tiempos en seg) para el MIDI clock variable. Cacheada.
    Dobla los golpes del click a negras usando el factor (click en corcheas = x2)."""
    d = _admin_audio_dir(numero)
    cache = d / "beatgrid.json"
    if cache.exists():
        try:
            return json.loads(cache.read_text(encoding="utf-8"))
        except Exception:
            pass
    res = _detectar_beats_click(numero)
    if not res.get("ok"):
        return {"ok": False, "grid": [], "bpm": 0}
    beats = res.get("beats", [])
    bpm_prom = res.get("bpm_prom", 0) or 0
    can = cargar_biblioteca().get(numero) or {}
    try:
        song_bpm = float(can.get("tempo") or 0)
    except Exception:
        song_bpm = 0.0
    factor = 1
    if song_bpm > 0 and bpm_prom > 0:
        r = bpm_prom / song_bpm
        f = int(round(r))
        if 1 <= f <= 4 and abs(r - f) <= 0.2:
            factor = f
    grid = [round(float(beats[i]), 4) for i in range(0, len(beats), factor)]
    bpm_real = int(round(bpm_prom / factor)) if factor else int(bpm_prom)
    out = {"ok": True, "grid": grid, "bpm": bpm_real}
    try:
        d.mkdir(parents=True, exist_ok=True)
        cache.write_text(json.dumps(out), encoding="utf-8")
    except Exception:
        pass
    return out


@app.route("/api/live/beatgrid/<int:numero>")
def api_live_beatgrid(numero):
    """NeuralPlay: grilla de negras para el MIDI clock (sigue cambios de tempo)."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    try:
        return jsonify(_beatgrid_musical(numero))
    except Exception as e:
        logging.error("beatgrid %s: %s", numero, e)
        return jsonify({"ok": False, "grid": [], "bpm": 0})


@app.route("/api/live/render/<int:numero>/<t>", methods=["POST"])
def api_live_render(numero, t):
    """NeuralPlay: dispara el render de un tono (semitonos). Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    try:
        n = int(t)
    except ValueError:
        return jsonify({"listo": False, "error": "tono invalido"}), 400
    if _tono_listo(numero, n):
        return jsonify({"listo": True})
    if not _stems_originales(numero):
        return jsonify({"listo": False, "error": "sin pistas"})
    d = _carpeta_tono(numero, n)
    if not (d.exists() and (d / ".lock").exists()):
        threading.Thread(target=_render_tono, args=(numero, n, org_actual()), daemon=True).start()
    return jsonify({"listo": False, "estado": "procesando"})


@app.route("/api/live/render/<int:numero>/<t>/estado")
def api_live_render_estado(numero, t):
    """NeuralPlay: progreso del render de un tono. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    try:
        n = int(t)
    except ValueError:
        return jsonify({"listo": False})
    if _tono_listo(numero, n):
        return jsonify({"listo": True})
    prog = ""
    lock = _carpeta_tono(numero, n) / ".lock"
    if lock.exists():
        try:
            prog = lock.read_text().strip()
        except Exception:
            prog = ""
    return jsonify({"listo": False, "progreso": prog})


@app.route("/api/live/setlist/crear", methods=["POST"])
def api_live_setlist_crear():
    if not _token_ok(): return jsonify({"ok": False, "error": "unauthorized"}), 403
    nombre = (request.values.get("nombre") or "").strip() or "Sin nombre"
    fecha = (request.values.get("fecha") or "").strip() or _hoy_iso()
    cfg = get_config()
    nuevo = {"id": secrets.token_hex(6), "nombre": nombre[:80], "fecha": fecha, "canciones": [], "creado": _ahora_iso()}
    cfg.setdefault("setlists", []).append(nuevo)
    guardar_config(cfg)
    return jsonify({"ok": True, "id": nuevo["id"]})


@app.route("/api/live/setlist/<sid>/eliminar", methods=["POST"])
def api_live_setlist_eliminar(sid):
    if not _token_ok(): return jsonify({"ok": False}), 403
    cfg = get_config()
    cfg["setlists"] = [x for x in cfg.get("setlists", []) if x["id"] != sid]
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/renombrar", methods=["POST"])
def api_live_setlist_renombrar(sid):
    if not _token_ok(): return jsonify({"ok": False}), 403
    nombre = (request.values.get("nombre") or "").strip()
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False}), 404
    if nombre: s["nombre"] = nombre[:80]
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/agregar", methods=["POST"])
def api_live_setlist_agregar(sid):
    if not _token_ok(): return jsonify({"ok": False}), 403
    try: numero = int(request.values.get("numero", ""))
    except ValueError: return jsonify({"ok": False, "error": "numero"}), 400
    tono = (request.values.get("tono") or "").strip() or None
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False, "error": "setlist"}), 404
    if numero not in cargar_biblioteca(): return jsonify({"ok": False, "error": "cancion"}), 404
    nuevo = {"id": numero, "tono": tono}
    cs = s.setdefault("canciones", [])
    after = request.values.get("after")
    pos = None
    if after:
        try: aid = int(after)
        except ValueError: aid = None
        if aid is not None:
            pos = next((i for i, c in enumerate(cs) if c.get("id") == aid), None)
    if pos is not None:
        cs.insert(pos + 1, nuevo)
    else:
        cs.append(nuevo)
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/quitar/<int:numero>", methods=["POST"])
def api_live_setlist_quitar(sid, numero):
    if not _token_ok(): return jsonify({"ok": False}), 403
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if s:
        s["canciones"] = [c for c in s.get("canciones", []) if c.get("id") != numero]
        guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/tono/<int:numero>", methods=["POST"])
def api_live_setlist_tono(sid, numero):
    if not _token_ok(): return jsonify({"ok": False}), 403
    tono = (request.values.get("tono") or "").strip()
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False}), 404
    item = next((c for c in s.get("canciones", []) if c.get("id") == numero), None)
    if not item: return jsonify({"ok": False}), 404
    base = cargar_biblioteca().get(numero, {})
    item["tono"] = None if (not tono or tono == base.get("tono", "")) else tono
    guardar_config(cfg)
    return jsonify({"ok": True, "tono": item["tono"]})


@app.route("/api/live/setlist/<sid>/mover", methods=["POST"])
def api_live_setlist_mover(sid):
    if not _token_ok(): return jsonify({"ok": False}), 403
    try: numero = int(request.values.get("numero", 0))
    except ValueError: numero = 0
    direccion = request.values.get("direccion", "")
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False}), 404
    cs = s.get("canciones", [])
    idx = next((i for i, c in enumerate(cs) if c.get("id") == numero), -1)
    if idx >= 0:
        j = idx - 1 if direccion == "arriba" else idx + 1
        if 0 <= j < len(cs):
            cs[idx], cs[j] = cs[j], cs[idx]
            guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/tonos/<int:numero>")
def api_live_tonos(numero):
    if not _token_ok(): return jsonify({"ok": False}), 403
    base = cargar_biblioteca().get(numero, {})
    rendered = []
    if _stems_originales(numero): rendered.append(0)
    d = dir_pistas(_cur_org()) / str(numero)
    if d.is_dir():
        for sub in d.iterdir():
            if sub.is_dir() and sub.name.startswith("tono_"):
                try: n = int(sub.name[5:])
                except ValueError: continue
                if _tono_listo(numero, n): rendered.append(n)
    rendered = sorted(set(rendered))
    orig = base.get("tono", "")
    from transposicion import NOTAS_BEMOL, transponer_acorde, usar_sostenidos, NOTA_A_INDICE
    import re as _re
    def _raiz(t):
        m=_re.match(r"^([A-G][#b]?)", t or ""); return m.group(1) if m else None
    keys=[]
    ro=_raiz(orig)
    for k in NOTAS_BEMOL:
        if ro and ro in NOTA_A_INDICE and k in NOTA_A_INDICE:
            d=(NOTA_A_INDICE[k]-NOTA_A_INDICE[ro])%12
            if d>6: d-=12
        else:
            d=0
        ts=transponer_acorde(orig,d,usar_sost=True) if orig else k
        nombre=transponer_acorde(orig,d,usar_sost=usar_sostenidos(ts)) if orig else k
        keys.append({"nombre":nombre,"semitonos":d,"rendered":d in rendered})
    return jsonify({"ok": True, "numero": numero, "tono_original": orig, "rendered": rendered, "keys": keys})


@app.route("/api/live/mix/<int:numero>", methods=["GET"])
def api_live_mix_get(numero):
    if not _token_ok(): return jsonify({"ok": False}), 403
    f = dir_pistas(_cur_org()) / str(numero) / "mix.json"
    if f.is_file():
        try:
            return jsonify({"ok": True, "mix": json.loads(f.read_text())})
        except Exception:
            pass
    return jsonify({"ok": True, "mix": None})


@app.route("/api/live/mix/<int:numero>", methods=["POST"])
def api_live_mix_save(numero):
    if not _token_ok(): return jsonify({"ok": False}), 403
    try:
        mix = json.loads(request.values.get("data", ""))
    except Exception:
        return jsonify({"ok": False, "error": "json"}), 400
    d = dir_pistas(_cur_org()) / str(numero)
    d.mkdir(parents=True, exist_ok=True)
    (d / "mix.json").write_text(json.dumps(mix, ensure_ascii=False))
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/orden", methods=["POST"])
def api_live_setlist_orden(sid):
    if not _token_ok(): return jsonify({"ok": False}), 403
    try:
        order = [int(x) for x in request.values.get("ids", "").split(",") if x.strip()]
    except ValueError:
        return jsonify({"ok": False}), 400
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False}), 404
    cs = s.get("canciones", [])
    buckets = {}
    for c in cs:
        buckets.setdefault(c.get("id"), []).append(c)
    newcs = []
    for i in order:
        if buckets.get(i):
            newcs.append(buckets[i].pop(0))
    for c in cs:
        if c not in newcs:
            newcs.append(c)
    s["canciones"] = newcs
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live/setlist/<sid>/mix", methods=["POST"])
def api_live_setlist_mix(sid):
    """Guarda la mezcla (faders/mutes/buses/master) de las canciones de un repertorio.
    Recibe data = { "<songId>": {master,tracks,buses}, ... }. Token."""
    if not _token_ok(): return jsonify({"ok": False}), 403
    try:
        data = json.loads(request.values.get("data", ""))
    except Exception:
        return jsonify({"ok": False, "error": "json"}), 400
    cfg = get_config()
    s = next((x for x in cfg.get("setlists", []) if x["id"] == sid), None)
    if not s: return jsonify({"ok": False}), 404
    for c in s.get("canciones", []):
        k = str(c.get("id"))
        if k in data:
            c["mix"] = data[k]
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/cancion/<int:numero>")
@login_required("musico")
def api_cancion(numero):
    biblioteca = cargar_biblioteca()
    if numero not in biblioteca:
        return jsonify({"error": "No encontrada"}), 404
    return jsonify(biblioteca[numero])


# ───────────────────────── Rutas admin ─────────────────────────
@app.route("/api/version")
def api_version():
    """Ultima version publicada de NeuralSync (para el aviso de actualizacion)."""
    return jsonify({"ok": True, "version": get_config().get("neuralworship_version", "1.0")})


DESCARGAS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "descargas")
ARCHIVOS_DESCARGA = {
    "mac": "NeuralSync-Mac.zip",
    "windows": "NeuralSync-Windows-Setup.exe",
}
ARCHIVOS_DESCARGA_NP = {
    "mac": "NeuralPlay-Mac.zip",
    "windows": "NeuralPlay-Windows-Setup.exe",
}

@app.route("/admin/descargar/<plataforma>")
@login_required("admin")
def admin_descargar(plataforma):
    tabla = ARCHIVOS_DESCARGA_NP if request.args.get("app") == "neuralplay" else ARCHIVOS_DESCARGA
    nombre = tabla.get(plataforma)
    if not nombre:
        abort(404)
    ruta = os.path.join(DESCARGAS_DIR, nombre)
    if not os.path.exists(ruta):
        flash("Ese instalador todavía no está disponible.", "error")
        return redirect(url_for("admin") + "#live")
    return send_file(ruta, as_attachment=True, download_name=nombre)


@app.route("/admin")
@login_required("admin")
def admin():
    biblioteca = cargar_biblioteca()
    biblioteca_ordenada = sorted(biblioteca.values(), key=lambda c: c["numero"])
    cfg = get_config()
    live_activo = is_live_activo(cfg)
    # Hace cuántos segundos fue el último heartbeat
    ultimo = cfg.get("ultimo_heartbeat")
    if ultimo:
        import time as _t
        segundos_atras = int(_t.time() - ultimo)
    else:
        segundos_atras = None
    # Setlists ordenados por fecha ASC
    setlists = sorted(
        cfg.get("setlists", []),
        key=lambda s: (s.get("fecha", "9999"), s.get("creado", ""))
    )
    # Para cada setlist, agregar info de las canciones
    biblioteca_dict = {c["numero"]: c for c in biblioteca_ordenada}
    setlists_full = []
    for s in setlists:
        canciones = []
        for item in s.get("canciones", []):
            num = item.get("id")
            if num in biblioteca_dict:
                base = biblioteca_dict[num]
                tono_efectivo = item.get("tono") or base.get("tono", "")
                canciones.append({
                    "id": num,
                    "titulo": base["titulo"],
                    "artista": base.get("artista", ""),
                    "tono": tono_efectivo,
                    "tono_original": base.get("tono", ""),
                    "es_override": bool(item.get("tono") and item["tono"] != base.get("tono", "")),
                })
        setlists_full.append({
            **s,
            "canciones_resolved": canciones,
        })
    pads_lista = []
    for _p in _pads_cargar():
        pads_lista.append({**_p, **_pad_estado(_p["id"])})
    return render_template(
        "admin.html",
        biblioteca=biblioteca_ordenada,
        pads=pads_lista,
        setlists=setlists_full,
        live_activo=live_activo,
        live_ip=cfg.get("mac_local_ip"),
        live_token=cfg.get("live_token", ""),
        ultimo_heartbeat_seg=segundos_atras,
        hoy=_hoy_iso(),
        usuarios_pendientes=usuarios.listar_usuarios(estado="pendiente", org_id=org_actual()),
        usuarios_activos=[u for u in usuarios.listar_usuarios(estado="activo", org_id=org_actual()) if u.get("rol") != "admin"],
        org_info=usuarios.obtener_organizacion(org_actual()),
        musicos_activos_n=usuarios.contar_musicos_activos(org_actual()),
        invitaciones_pendientes=usuarios.listar_invitaciones(org_actual(), "pendiente"),
        usuario_actual=get_usuario_actual(),
        email_configurado=emails_module.email_configurado(),
        desc_mac=os.path.exists(os.path.join(DESCARGAS_DIR, ARCHIVOS_DESCARGA["mac"])),
        desc_win=os.path.exists(os.path.join(DESCARGAS_DIR, ARCHIVOS_DESCARGA["windows"])),
        desc_np_mac=os.path.exists(os.path.join(DESCARGAS_DIR, ARCHIVOS_DESCARGA_NP["mac"])),
        desc_np_win=os.path.exists(os.path.join(DESCARGAS_DIR, ARCHIVOS_DESCARGA_NP["windows"])),
    )


@app.route("/admin/subir", methods=["POST"])
@login_required("admin")
def admin_subir():
    if "archivo" not in request.files:
        flash("No se seleccionó archivo", "error")
        return redirect(url_for("admin"))
    archivo = request.files["archivo"]
    if archivo.filename == "":
        flash("Archivo vacío", "error")
        return redirect(url_for("admin"))
    if not archivo.filename.endswith(".json"):
        flash("Solo se aceptan archivos .json", "error")
        return redirect(url_for("admin"))
    try:
        contenido = archivo.read()
        datos = json.loads(contenido.decode("utf-8"))
        if "numero" not in datos or "titulo" not in datos or "secciones" not in datos:
            flash("El JSON no tiene los campos requeridos (numero, titulo, secciones)", "error")
            return redirect(url_for("admin"))
    except json.JSONDecodeError:
        flash("Archivo no es JSON válido", "error")
        return redirect(url_for("admin"))
    nombre_seguro = secure_filename(archivo.filename)
    destino = dir_canciones(org_actual()) / nombre_seguro
    with open(destino, "wb") as f:
        f.write(contenido)
    flash(f"✓ Canción '{datos['titulo']}' (#{datos['numero']}) subida correctamente", "success")
    return redirect(url_for("admin"))


@app.route("/admin/eliminar/<int:numero>", methods=["POST"])
@login_required("admin")
def admin_eliminar(numero):
    biblioteca = cargar_biblioteca()
    if numero not in biblioteca:
        flash("Canción no encontrada", "error")
        return redirect(url_for("admin"))
    for archivo in dir_canciones(org_actual()).glob("*.json"):
        try:
            with open(archivo, "r", encoding="utf-8") as f:
                datos = json.load(f)
            if datos.get("numero") == numero:
                archivo.unlink()
                carpeta_p = dir_pistas(_cur_org()) / str(numero)
                if carpeta_p.is_dir():
                    try:
                        shutil.rmtree(str(carpeta_p))
                    except Exception:
                        pass
                flash(f"✓ Canción #{numero} eliminada", "success")
                # Removerla de TODOS los setlists si estaba
                cfg = get_config()
                for s in cfg.get("setlists", []):
                    s["canciones"] = [c for c in s.get("canciones", []) if c.get("id") != numero]
                guardar_config(cfg)
                break
        except Exception:
            continue
    return redirect(url_for("admin"))


# ───────── CRUD de setlists ─────────
@app.route("/admin/setlist/crear", methods=["POST"])
@login_required("admin")
def admin_setlist_crear():
    nombre = request.form.get("nombre", "").strip() or "Sin nombre"
    fecha = request.form.get("fecha", "").strip() or _hoy_iso()
    cfg = get_config()
    nuevo = {
        "id": secrets.token_hex(6),
        "nombre": nombre[:80],
        "fecha": fecha,
        "canciones": [],
        "creado": _ahora_iso(),
    }
    cfg.setdefault("setlists", []).append(nuevo)
    guardar_config(cfg)
    flash(f"✓ Setlist '{nombre}' creado", "success")
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/editar", methods=["POST"])
@login_required("admin")
def admin_setlist_editar(setlist_id):
    nombre = request.form.get("nombre", "").strip()
    fecha = request.form.get("fecha", "").strip()
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if not s:
        flash("Setlist no encontrado", "error")
        return redirect(url_for("admin"))
    if nombre:
        s["nombre"] = nombre[:80]
    if fecha:
        s["fecha"] = fecha
    guardar_config(cfg)
    flash("✓ Setlist actualizado", "success")
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/eliminar", methods=["POST"])
@login_required("admin")
def admin_setlist_eliminar(setlist_id):
    cfg = get_config()
    antes = len(cfg.get("setlists", []))
    cfg["setlists"] = [s for s in cfg.get("setlists", []) if s["id"] != setlist_id]
    if len(cfg["setlists"]) < antes:
        guardar_config(cfg)
        flash("✓ Setlist eliminado", "success")
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/agregar", methods=["POST"])
@login_required("admin")
def admin_setlist_agregar(setlist_id):
    numero = request.form.get("numero", "").strip()
    try:
        numero = int(numero)
    except ValueError:
        flash("Número inválido", "error")
        return redirect(url_for("admin"))
    biblioteca = cargar_biblioteca()
    if numero not in biblioteca:
        flash("Canción no encontrada", "error")
        return redirect(url_for("admin"))
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if not s:
        flash("Setlist no encontrado", "error")
        return redirect(url_for("admin"))
    # Evitar duplicados
    if any(c.get("id") == numero for c in s.get("canciones", [])):
        flash("La canción ya está en el setlist", "error")
        return redirect(url_for("admin"))
    s.setdefault("canciones", []).append({"id": numero, "tono": None})
    guardar_config(cfg)
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/quitar/<int:numero>", methods=["POST"])
@login_required("admin")
def admin_setlist_quitar(setlist_id, numero):
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if s:
        s["canciones"] = [c for c in s.get("canciones", []) if c.get("id") != numero]
        guardar_config(cfg)
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/mover", methods=["POST"])
@login_required("admin")
def admin_setlist_mover(setlist_id):
    numero = int(request.form.get("numero", 0))
    direccion = request.form.get("direccion", "")
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if not s:
        return redirect(url_for("admin"))
    canciones = s.get("canciones", [])
    idx = next((i for i, c in enumerate(canciones) if c.get("id") == numero), -1)
    if idx < 0:
        return redirect(url_for("admin"))
    if direccion == "arriba" and idx > 0:
        canciones[idx-1], canciones[idx] = canciones[idx], canciones[idx-1]
    elif direccion == "abajo" and idx < len(canciones) - 1:
        canciones[idx+1], canciones[idx] = canciones[idx], canciones[idx+1]
    guardar_config(cfg)
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/transponer/<int:numero>", methods=["POST"])
@login_required("admin")
def admin_setlist_transponer(setlist_id, numero):
    """Sube/baja N semitonos el tono de una canción en un setlist."""
    delta = int(request.form.get("delta", 0))
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if not s:
        return redirect(url_for("admin"))
    item = next((c for c in s.get("canciones", []) if c.get("id") == numero), None)
    if not item:
        return redirect(url_for("admin"))
    biblioteca = cargar_biblioteca()
    base = biblioteca.get(numero, {})
    tono_actual = item.get("tono") or base.get("tono", "")
    tono_abs = (request.form.get("tono") or "").strip()
    if tono_abs:
        item["tono"] = None if tono_abs == base.get("tono", "") else tono_abs
        guardar_config(cfg)
        return redirect(url_for("admin"))
    if not tono_actual:
        return redirect(url_for("admin"))
    # Calcular nuevo tono (smart: bemoles/sostenidos)
    from transposicion import transponer_acorde, usar_sostenidos
    tono_sost = transponer_acorde(tono_actual, delta, usar_sost=True)
    usar_sost = usar_sostenidos(tono_sost)
    nuevo_tono = transponer_acorde(tono_actual, delta, usar_sost=usar_sost)
    # Si vuelve al original, limpiar override
    if nuevo_tono == base.get("tono", ""):
        item["tono"] = None
    else:
        item["tono"] = nuevo_tono
    guardar_config(cfg)
    return redirect(url_for("admin"))


@app.route("/admin/setlist/<setlist_id>/resetear/<int:numero>", methods=["POST"])
@login_required("admin")
def admin_setlist_resetear_tono(setlist_id, numero):
    cfg = get_config()
    s = next((s for s in cfg.get("setlists", []) if s["id"] == setlist_id), None)
    if s:
        item = next((c for c in s.get("canciones", []) if c.get("id") == numero), None)
        if item:
            item["tono"] = None
            guardar_config(cfg)
    return redirect(url_for("admin"))


@app.route("/admin/cambiar_password", methods=["POST"])
@login_required("admin")
def admin_cambiar_password():
    cual = request.form.get("cual", "")
    nueva = request.form.get("nueva", "")
    if not nueva or len(nueva) < 6:
        flash("Contraseña debe tener al menos 6 caracteres", "error")
        return redirect(url_for("admin"))
    cfg = get_config()
    if cual == "musicos":
        cfg["password_musicos"] = hash_password(nueva)
        flash("✓ Password de músicos actualizado", "success")
    elif cual == "admin":
        cfg["password_admin"] = hash_password(nueva)
        flash("✓ Password de admin actualizado", "success")
    else:
        flash("Tipo de password inválido", "error")
        return redirect(url_for("admin"))
    guardar_config(cfg)
    return redirect(url_for("admin"))


# ───────────────────────── Live (Mac local) ─────────────────────────
# Timeout: si pasaron más de 90 segundos sin heartbeat, el live se considera apagado
HEARTBEAT_TIMEOUT_SEG = 90


def is_live_activo(cfg=None):
    """True si el Mac mandó heartbeat recientemente."""
    if cfg is None:
        cfg = get_config()
    if not cfg.get("live_activo"):
        return False
    ultimo = cfg.get("ultimo_heartbeat")
    if not ultimo:
        return False
    import time as _t
    return (_t.time() - ultimo) < HEARTBEAT_TIMEOUT_SEG


@app.route("/api/live_ping", methods=["POST"])
def api_live_ping():
    """El puente.py del Mac llama esto cada 30s para avisar que está vivo."""
    data = request.get_json(silent=True) or {}
    cfg = get_config()
    # Validar token
    if data.get("token") != cfg.get("live_token"):
        return jsonify({"ok": False, "error": "invalid token"}), 403
    ip = data.get("ip", "").strip()
    if not ip:
        return jsonify({"ok": False, "error": "missing ip"}), 400
    accion = data.get("accion", "ping")  # "ping" o "bye"
    import time as _t
    if accion == "bye":
        cfg["live_activo"] = False
        cfg["ultimo_heartbeat"] = None
        print(f"📴 Mac local desconectado (bye explícito)")
    else:
        cfg["live_activo"] = True
        cfg["mac_local_ip"] = ip
        cfg["ultimo_heartbeat"] = _t.time()
    guardar_config(cfg)
    return jsonify({"ok": True})


@app.route("/api/live_status")
def api_live_status():
    """El frontend lo consulta para saber si hay live activo."""
    cfg = get_config()
    activo = is_live_activo(cfg)
    return jsonify({
        "activo": activo,
        "ip": cfg.get("mac_local_ip") if activo else None,
    })


@app.route("/admin/live_off", methods=["POST"])
@login_required("admin")
def admin_live_off():
    """Forzar apagado del live desde el panel admin."""
    cfg = get_config()
    cfg["live_activo"] = False
    cfg["ultimo_heartbeat"] = None
    guardar_config(cfg)
    flash("✓ Live forzado a OFF", "success")
    return redirect(url_for("admin"))


# ───────────────────────── Gestión de usuarios (admin) ─────────────────────────
@app.route("/admin/usuario/<int:user_id>/aprobar", methods=["POST"])
@login_required("admin")
def admin_usuario_aprobar(user_id):
    org = org_actual()
    u = usuarios.buscar_por_id(user_id)
    if not u or u.get("org_id") != org:
        flash("Usuario no encontrado", "error")
        return redirect(url_for("admin"))
    org_obj = usuarios.obtener_organizacion(org)
    if (u.get("rol") == "musico" and org_obj
            and usuarios.contar_musicos_activos(org) >= int(org_obj.get("max_musicos") or 0)):
        flash("Alcanzaste el límite de asientos de tu paquete. Ampliá el plan o quitá un músico para aprobar a otro.", "error")
        return redirect(url_for("admin"))
    if usuarios.aprobar_usuario(user_id, org):
        # Mandar email de bienvenida
        nombre_completo = f"{u['nombre']} {u['apellido']}"
        ok, _ = emails_module.enviar_email_bienvenida(u["email"], nombre_completo)
        if ok:
            flash(f"✓ {nombre_completo} aprobado y notificado por email", "success")
        else:
            flash(f"✓ {nombre_completo} aprobado (no se pudo mandar el email)", "success")
    else:
        flash("No se pudo aprobar (¿ya estaba aprobado?)", "error")
    return redirect(url_for("admin"))


@app.route("/admin/usuario/<int:user_id>/eliminar", methods=["POST"])
@login_required("admin")
def admin_usuario_eliminar(user_id):
    org = org_actual()
    u = usuarios.buscar_por_id(user_id)
    if not u or u.get("org_id") != org:
        flash("Usuario no encontrado", "error")
        return redirect(url_for("admin"))
    # No permitir borrar al dueño de la organización
    org_obj = usuarios.obtener_organizacion(org)
    if org_obj and org_obj.get("owner_user_id") == user_id:
        flash("No podés borrar al dueño de la organización", "error")
        return redirect(url_for("admin"))
    # No permitir auto-eliminación
    if session.get("user_id") == user_id:
        flash("No podés borrar tu propia cuenta", "error")
        return redirect(url_for("admin"))
    if usuarios.eliminar_usuario(user_id, org):
        flash(f"✓ Usuario {u['nombre']} {u['apellido']} eliminado", "success")
    return redirect(url_for("admin"))


@app.route("/admin/usuario/<int:user_id>/reset_password", methods=["POST"])
@login_required("admin")
def admin_usuario_reset_password(user_id):
    """Admin manda manualmente un código de reset (por si el usuario no recibe el email original)."""
    u = usuarios.buscar_por_id(user_id)
    if not u or u.get("org_id") != org_actual():
        flash("Usuario no encontrado", "error")
        return redirect(url_for("admin"))
    codigo = usuarios.crear_codigo_reset(user_id)
    nombre_completo = f"{u['nombre']} {u['apellido']}"
    ok, msg = emails_module.enviar_email_codigo_reset(u["email"], nombre_completo, codigo)
    if ok:
        flash(f"✓ Código de reset enviado a {u['email']}", "success")
    else:
        flash(f"⚠️ No se pudo mandar email: {msg}. Código: {codigo}", "error")
    return redirect(url_for("admin"))


@app.route("/admin/org/nombre", methods=["POST"])
@login_required("admin")
def admin_org_nombre():
    nombre = request.form.get("org_nombre", "").strip()
    if not nombre:
        flash("El nombre no puede estar vacío", "error")
        return redirect(url_for("admin"))
    cfg = get_config()
    cfg["org_nombre"] = nombre[:80]
    guardar_config(cfg)
    flash("✓ Nombre de la organización actualizado", "success")
    return redirect(url_for("admin"))


@app.route("/admin/org/logo", methods=["POST"])
@login_required("admin")
def admin_org_logo():
    f = request.files.get("logo")
    if not f or not f.filename:
        flash("Elegí una imagen", "error")
        return redirect(url_for("admin"))
    ext = f.filename.rsplit(".", 1)[-1].lower() if "." in f.filename else ""
    if ext not in ("png", "jpg", "jpeg", "webp"):
        flash("Formato no válido (usá PNG, JPG o WEBP)", "error")
        return redirect(url_for("admin"))
    try:
        f.save(str(BASE_DIR / "static" / "logo.png"))
        flash("✓ Logo actualizado", "success")
    except Exception as e:
        flash("No se pudo guardar el logo: %s" % e, "error")
    return redirect(url_for("admin"))


@app.route("/admin/perfil", methods=["POST"])
@login_required("admin")
def admin_perfil():
    uid = session.get("user_id")
    if not uid:
        flash("No hay una cuenta asociada a esta sesión (acceso de emergencia).", "error")
        return redirect(url_for("admin"))
    ok, msg = usuarios.actualizar_perfil(
        uid,
        request.form.get("nombre", ""),
        request.form.get("apellido", ""),
        request.form.get("email", ""),
    )
    if ok:
        session["nombre"] = request.form.get("nombre", "").strip()
        flash("✓ Perfil actualizado", "success")
    else:
        flash(msg, "error")
    return redirect(url_for("admin"))


@app.route("/admin/mi_password", methods=["POST"])
@login_required("admin")
def admin_cambiar_mi_password():
    """El admin en sesión cambia su propia contraseña de cuenta."""
    uid = session.get("user_id")
    if not uid:
        flash("No hay una cuenta asociada a esta sesión (acceso de emergencia).", "error")
        return redirect(url_for("admin"))
    nueva = request.form.get("nueva", "")
    if not nueva or len(nueva) < 6:
        flash("La contraseña debe tener al menos 6 caracteres", "error")
        return redirect(url_for("admin"))
    ok, msg = usuarios.cambiar_password_directo(uid, nueva)
    if ok:
        flash("✓ Tu contraseña fue actualizada", "success")
    else:
        flash(msg, "error")
    return redirect(url_for("admin"))


@app.route("/admin/crear_admin", methods=["POST"])
@login_required("admin")
def admin_crear_admin():
    """Crear un nuevo admin desde el panel (solo otros admins pueden)."""
    nombre = request.form.get("nombre", "").strip()
    apellido = request.form.get("apellido", "").strip()
    email = request.form.get("email", "").strip().lower()
    password = request.form.get("password", "")
    ok, resultado = usuarios.crear_usuario(nombre, apellido, email, password, rol="admin", estado="activo", org_id=org_actual())
    if ok:
        flash(f"✓ Admin {nombre} {apellido} creado", "success")
    else:
        flash(f"Error: {resultado}", "error")
    return redirect(url_for("admin"))


@app.route("/admin/invitar", methods=["POST"])
@login_required("admin")
def admin_invitar():
    """El admin invita a un músico por email. Le llega un link para entrar directo."""
    org = org_actual()
    email = request.form.get("email", "").strip().lower()
    if not email:
        flash("Ingresá un email", "error")
        return redirect(url_for("admin"))
    org_obj = usuarios.obtener_organizacion(org)
    if org_obj and usuarios.contar_musicos_activos(org) >= int(org_obj.get("max_musicos") or 0):
        flash("Alcanzaste el límite de asientos de tu paquete. Ampliá el plan para invitar a más músicos.", "error")
        return redirect(url_for("admin"))
    ok, res = usuarios.crear_invitacion(org, email)
    if not ok:
        flash(res, "error")
        return redirect(url_for("admin"))
    link = url_for("aceptar_invitacion", token=res, _external=True)
    org_nombre = (org_obj or {}).get("nombre", "tu organización")
    okmail, _ = emails_module.enviar_email_invitacion(email, org_nombre, link)
    if okmail:
        flash("✓ Invitación enviada a %s" % email, "success")
    else:
        flash("Invitación creada, pero no se pudo enviar el correo. Enlace: %s" % link, "error")
    return redirect(url_for("admin"))


@app.route("/admin/invitacion/<int:inv_id>/cancelar", methods=["POST"])
@login_required("admin")
def admin_invitacion_cancelar(inv_id):
    usuarios.eliminar_invitacion(inv_id, org_actual())
    flash("Invitación cancelada", "success")
    return redirect(url_for("admin"))


@app.route("/invitacion/<token>", methods=["GET", "POST"])
def aceptar_invitacion(token):
    """El músico invitado crea su cuenta (queda ACTIVA al instante) y entra."""
    inv = usuarios.obtener_invitacion(token)
    if not inv:
        flash("La invitación no es válida o ya venció.", "error")
        return redirect(url_for("login"))
    org = usuarios.obtener_organizacion(inv["org_id"])
    if request.method == "POST":
        nombre = request.form.get("nombre", "").strip()
        apellido = request.form.get("apellido", "").strip()
        password = request.form.get("password", "")
        password2 = request.form.get("password2", "")
        prev = dict(nombre=nombre, apellido=apellido)
        if not nombre or not apellido or not password:
            flash("Completá todos los campos", "error")
            return render_template("invitacion.html", inv=inv, org=org, **prev)
        if password != password2:
            flash("Las contraseñas no coinciden", "error")
            return render_template("invitacion.html", inv=inv, org=org, **prev)
        if org and usuarios.contar_musicos_activos(inv["org_id"]) >= int(org.get("max_musicos") or 0):
            flash("La organización alcanzó su límite de asientos. Pedile al administrador que amplíe el plan.", "error")
            return render_template("invitacion.html", inv=inv, org=org, **prev)
        ok, res = usuarios.crear_usuario(nombre, apellido, inv["email"], password,
                                         rol="musico", estado="activo", org_id=inv["org_id"])
        if not ok:
            flash(res, "error")
            return render_template("invitacion.html", inv=inv, org=org, **prev)
        usuarios.marcar_invitacion_usada(token)
        session.permanent = True
        session.clear()
        session["user_id"] = res
        session["rol"] = "musico"
        session["nombre"] = nombre
        session["org_id"] = inv["org_id"]
        flash("✓ ¡Bienvenido a %s!" % (org["nombre"] if org else "NeuralWorship"), "success")
        return redirect(url_for("principal"))
    return render_template("invitacion.html", inv=inv, org=org)


# ───────────────────────── Main ─────────────────────────
# ---- Reproductor de practica (stems / modo ensayo) ----
CARPETA_PISTAS = BASE_DIR / "pistas"
CARPETA_PISTAS.mkdir(exist_ok=True)
_EXT_AUDIO_ENSAYO = (".mp3", ".m4a", ".ogg", ".wav")
FAMILIAS_FIJAS = {"Batería", "Percusión", "Guía", "Click"}


def _carpeta_tono(numero, n):
    if n == 0:
        return dir_pistas(_cur_org()) / str(numero)
    return dir_pistas(_cur_org()) / str(numero) / ("tono_" + str(n))


def _stems_originales(numero):
    d = dir_pistas(_cur_org()) / str(numero)
    if not d.is_dir():
        return []
    return sorted([f.name for f in d.iterdir() if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO])


def _tono_listo(numero, n):
    if n == 0:
        return len(_stems_originales(numero)) > 0
    d = _carpeta_tono(numero, n)
    if not d.is_dir() or (d / ".lock").exists():
        return False
    orig = _stems_originales(numero)
    if not orig:
        return False
    hechos = set(f.name for f in d.iterdir() if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO)
    return all(o in hechos for o in orig)


def _carpeta_stems(numero, n):
    """Carpeta con los stems: originales (pistas/N/) si n==0, si no la del tono."""
    if n == 0:
        return dir_pistas(_cur_org()) / str(numero)
    return _carpeta_tono(numero, n)


def _tono_en_proceso(numero, n):
    """True si el tono se esta renderizando ahora mismo (existe .lock)."""
    if n == 0:
        return False
    d = _carpeta_tono(numero, n)
    return d.is_dir() and (d / ".lock").exists()


def _nombre_base_export(cancion, numero):
    """Solo el titulo para nombrar los .mid."""
    titulo = (cancion.get("titulo") or "").strip()
    base = secure_filename(titulo) if titulo else ""
    return base or ("cancion_%d" % numero)


def _invalidar_tonos(numero):
    """Borra los caches de tonos transpuestos: quedan inconsistentes al cambiar los stems."""
    base = dir_pistas(_cur_org()) / str(numero)
    if not base.is_dir():
        return
    for d in base.glob("tono_*"):
        if d.is_dir():
            try:
                shutil.rmtree(str(d))
            except Exception:
                pass
    adm = base / "_admin"
    if adm.is_dir():
        try:
            shutil.rmtree(str(adm))
        except Exception:
            pass


def _asegurar_web(numero, n, org=None):
    """Genera los proxys MP3 128k (modo ensayo) que falten para el tono n."""
    if org is not None:
        _set_org_ctx(org)
    d = _carpeta_tono(numero, n)
    web = d / "web"
    web.mkdir(parents=True, exist_ok=True)
    lock = web / ".lock"
    if lock.exists():
        return
    try:
        reales = [f for f in sorted(d.iterdir())
                  if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO]
        lock.write_text("0/" + str(len(reales)))
        hechos = 0
        for f in reales:
            out = web / (f.stem + ".mp3")
            if not out.exists():
                try:
                    subprocess.run(
                        ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y", "-i", str(f),
                         "-b:a", "128k", str(out)],
                        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=300)
                except Exception as e:
                    logging.error("proxy %s/%s %s: %s", numero, n, f.name, e)
            hechos += 1
            try:
                lock.write_text(str(hechos) + "/" + str(len(reales)))
            except Exception:
                pass
    finally:
        try:
            lock.unlink()
        except Exception:
            pass


def _web_listo(numero, n):
    """True si todos los stems reales del tono n tienen su proxy MP3."""
    if not _tono_listo(numero, n):
        return False
    d = _carpeta_tono(numero, n)
    web = d / "web"
    reales = [f.stem for f in d.iterdir()
              if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO]
    if not reales:
        return False
    return all((web / (s + ".mp3")).exists() for s in reales)


def _render_tono(numero, n, org=None):
    if org is not None:
        _set_org_ctx(org)
    d = _carpeta_tono(numero, n)
    d.mkdir(parents=True, exist_ok=True)
    lock = d / ".lock"
    orig = _stems_originales(numero)
    ratio = 2 ** (n / 12.0)
    try:
        lock.write_text("0/" + str(len(orig)))
        hechos = 0
        fam_map = _leer_familias(numero)
        for nombre in orig:
            entrada = dir_pistas(_cur_org()) / str(numero) / nombre
            salida = d / nombre
            if not salida.exists():
                fam = fam_map.get(nombre) or _familia_auto(Path(nombre).stem)
                if fam in FAMILIAS_FIJAS:
                    shutil.copy2(str(entrada), str(salida))
                else:
                    subprocess.run(
                        ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y", "-i", str(entrada),
                         "-af", "rubberband=pitch=" + repr(ratio), "-b:a", "192k", str(salida)],
                        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=300)
            hechos += 1
            try:
                lock.write_text(str(hechos) + "/" + str(len(orig)))
            except Exception:
                pass
    except Exception as e:
        logging.error("render tono %s/%s: %s", numero, n, e)
    finally:
        try:
            lock.unlink()
        except Exception:
            pass
    # proxys DESPUES de soltar el lock real: el reproductor ya ve el tono listo y baja los reales
    try:
        _asegurar_web(numero, n)
    except Exception as e:
        logging.error("proxys tono %s/%s: %s", numero, n, e)


# ───────────────────────── Pads ambientales ─────────────────────────
# Biblioteca GLOBAL de pads por tono (12 raíces), reutilizable por todas
# las canciones. Se sube UN pad base y el servidor genera los otros 11 por
# pitch-shift (rubberband), igual que los tonos de las canciones. Cada pack
# es un timbre distinto (Warm, Cinematic, etc.). NeuralPlay elige el pad
# según la raíz del tono de cada canción.
CARPETA_PADS = BASE_DIR / "pads"
CARPETA_PADS.mkdir(exist_ok=True)
_PADS_INDEX = CARPETA_PADS / "packs.json"
_EXT_AUDIO_PAD = (".wav", ".mp3", ".m4a", ".ogg", ".flac", ".aif", ".aiff")
PAD_ROOTS = ['C', 'C#', 'D', 'D#', 'E', 'F', 'F#', 'G', 'G#', 'A', 'A#', 'B']
_pads_lock = threading.Lock()


def _pad_root_idx(tono):
    """Raíz de un tono ('Db', 'Am', 'F#m'...) -> clase de altura 0-11, o None."""
    m = re.match(r'^\s*([A-G][#b]?)', tono or '')
    if not m:
        return None
    return NOTA_A_INDICE.get(m.group(1))


def _pads_cargar():
    try:
        return json.loads(_PADS_INDEX.read_text())
    except Exception:
        return []


def _pads_guardar(lst):
    try:
        _PADS_INDEX.write_text(json.dumps(lst, ensure_ascii=False, indent=2))
    except Exception as e:
        logging.error("guardar packs.json: %s", e)


def _pad_pack(pack_id):
    for p in _pads_cargar():
        if p.get("id") == pack_id:
            return p
    return None


def _pad_semitonos(base_idx, target_idx):
    """Semitonos mínimos (rango -5..+6) para ir de base_idx a target_idx."""
    d = (target_idx - base_idx) % 12
    if d > 6:
        d -= 12
    return d


def _pad_estado(pack_id):
    """Progreso/estado del render de un pack: {hechos,total,listo,render}."""
    d = CARPETA_PADS / pack_id
    total = 12
    lock = d / ".lock"
    if lock.exists():
        hechos = 0
        try:
            hechos, total = (int(x) for x in lock.read_text().split("/"))
        except Exception:
            pass
        return {"hechos": hechos, "total": total, "listo": False, "render": True}
    hechos = sum(1 for i in range(12) if (d / ("pad_%d.wav" % i)).is_file())
    return {"hechos": hechos, "total": total, "listo": hechos >= total, "render": False}


def _render_pad_pack(pack_id):
    """Genera los 12 tonos del pad desde el base, por pitch-shift (rubberband)."""
    pack = _pad_pack(pack_id)
    if not pack:
        return
    d = CARPETA_PADS / pack_id
    d.mkdir(parents=True, exist_ok=True)
    base = d / ("base" + pack.get("ext", ".wav"))
    if not base.is_file():
        logging.error("pad pack %s sin base", pack_id)
        return
    base_idx = int(pack.get("base_idx", 0))
    lock = d / ".lock"
    try:
        lock.write_text("0/12")
        hechos = 0
        for i in range(12):
            out = d / ("pad_%d.wav" % i)
            if not out.exists():
                semis = _pad_semitonos(base_idx, i)
                if semis == 0:
                    cmd = ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y",
                           "-i", str(base), str(out)]
                else:
                    ratio = 2 ** (semis / 12.0)
                    cmd = ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y",
                           "-i", str(base), "-af", "rubberband=pitch=" + repr(ratio),
                           str(out)]
                try:
                    subprocess.run(cmd, stdout=subprocess.DEVNULL,
                                   stderr=subprocess.DEVNULL, timeout=600)
                except Exception as e:
                    logging.error("render pad %s idx %d: %s", pack_id, i, e)
            hechos += 1
            try:
                lock.write_text(str(hechos) + "/12")
            except Exception:
                pass
    finally:
        try:
            lock.unlink()
        except Exception:
            pass


def _render_pad_async(pack_id):
    threading.Thread(target=_render_pad_pack, args=(pack_id,), daemon=True).start()


# ---- Admin: gestión de pads ----
@app.route("/admin/pads")
@login_required("admin")
def admin_pads():
    packs = []
    for p in _pads_cargar():
        packs.append({**p, **_pad_estado(p["id"])})
    return render_template("admin_pads.html", packs=packs, roots=PAD_ROOTS)


@app.route("/admin/pads/crear", methods=["POST"])
@login_required("admin")
def admin_pads_crear():
    nombre = (request.form.get("nombre") or "").strip()
    artista = (request.form.get("artista") or "").strip()
    base_root = (request.form.get("base_root") or "").strip()
    archivo = request.files.get("pad")
    if not nombre:
        flash("Ponele nombre al pack", "error")
        return redirect(url_for("admin_pads"))
    base_idx = _pad_root_idx(base_root)
    if base_idx is None:
        flash("Tono base inválido", "error")
        return redirect(url_for("admin_pads"))
    if not archivo or archivo.filename == "":
        flash("Subí el pad base", "error")
        return redirect(url_for("admin_pads"))
    ext = os.path.splitext(archivo.filename)[1].lower()
    if ext not in _EXT_AUDIO_PAD:
        flash("Formato no válido (wav/mp3/m4a/ogg/flac/aiff)", "error")
        return redirect(url_for("admin_pads"))
    with _pads_lock:
        lst = _pads_cargar()
        pack_id = secrets.token_hex(4)
        while any(x.get("id") == pack_id for x in lst):
            pack_id = secrets.token_hex(4)
        d = CARPETA_PADS / pack_id
        d.mkdir(parents=True, exist_ok=True)
        archivo.save(str(d / ("base" + ext)))
        portada_rel = ""
        pf = request.files.get("portada")
        if pf and pf.filename:
            pext = os.path.splitext(pf.filename)[1].lower()
            if pext in (".png", ".jpg", ".jpeg", ".webp"):
                try:
                    (BASE_DIR / "static" / "portadas").mkdir(parents=True, exist_ok=True)
                    pf.save(str(BASE_DIR / "static" / "portadas" / ("pad_" + pack_id + pext)))
                    portada_rel = "portadas/pad_" + pack_id + pext
                except Exception as e:
                    logging.error("portada pad %s: %s", pack_id, e)
        entrada = {"id": pack_id, "nombre": nombre, "artista": artista,
                   "base_root": PAD_ROOTS[base_idx], "base_idx": base_idx, "ext": ext}
        if portada_rel:
            import time as _t
            entrada["portada"] = portada_rel
            entrada["portada_ts"] = int(_t.time())
        lst.append(entrada)
        _pads_guardar(lst)
    _render_pad_async(pack_id)
    flash("Pad creado. Generando los 12 tonos…", "success")
    return redirect(url_for("admin_pads"))


@app.route("/admin/pads/<pack_id>/eliminar", methods=["POST"])
@login_required("admin")
def admin_pads_eliminar(pack_id):
    with _pads_lock:
        _pads_guardar([x for x in _pads_cargar() if x.get("id") != pack_id])
    d = CARPETA_PADS / pack_id
    if d.is_dir():
        try:
            shutil.rmtree(str(d))
        except Exception:
            pass
    for _ext in (".png", ".jpg", ".jpeg", ".webp"):
        _pf = BASE_DIR / "static" / "portadas" / ("pad_" + pack_id + _ext)
        try:
            if _pf.exists():
                _pf.unlink()
        except Exception:
            pass
    flash("Pad eliminado", "success")
    return redirect(url_for("admin_pads"))


@app.route("/admin/pads/<pack_id>/regenerar", methods=["POST"])
@login_required("admin")
def admin_pads_regenerar(pack_id):
    if not _pad_pack(pack_id):
        abort(404)
    d = CARPETA_PADS / pack_id
    for i in range(12):
        f = d / ("pad_%d.wav" % i)
        try:
            if f.exists():
                f.unlink()
        except Exception:
            pass
    _render_pad_async(pack_id)
    flash("Regenerando los 12 tonos…", "success")
    return redirect(url_for("admin_pads"))


@app.route("/admin/pads/<pack_id>/estado")
@login_required("admin")
def admin_pads_estado(pack_id):
    if not _pad_pack(pack_id):
        abort(404)
    return jsonify(_pad_estado(pack_id))


@app.route("/admin/pads/<pack_id>/editar", methods=["GET", "POST"])
@login_required("admin")
def admin_pads_editar(pack_id):
    pack = _pad_pack(pack_id)
    if not pack:
        abort(404)
    if request.method == "GET":
        return render_template("admin_pad_editar.html", p=pack, roots=PAD_ROOTS)
    nombre = (request.form.get("nombre") or "").strip()
    artista = (request.form.get("artista") or "").strip()
    base_root = (request.form.get("base_root") or "").strip()
    base_idx = _pad_root_idx(base_root)
    if not nombre or base_idx is None:
        flash("Datos invalidos", "error")
        return redirect(url_for("admin_pads_editar", pack_id=pack_id))
    d = CARPETA_PADS / pack_id
    regen = False
    with _pads_lock:
        lst = _pads_cargar()
        for it in lst:
            if it.get("id") != pack_id:
                continue
            it["nombre"] = nombre
            it["artista"] = artista
            if int(it.get("base_idx", 0)) != base_idx:
                it["base_idx"] = base_idx
                it["base_root"] = PAD_ROOTS[base_idx]
                regen = True
            nuevo = request.files.get("pad")
            if nuevo and nuevo.filename:
                ext = os.path.splitext(nuevo.filename)[1].lower()
                if ext in _EXT_AUDIO_PAD:
                    for old in d.glob("base.*"):
                        try:
                            old.unlink()
                        except Exception:
                            pass
                    d.mkdir(parents=True, exist_ok=True)
                    nuevo.save(str(d / ("base" + ext)))
                    it["ext"] = ext
                    regen = True
            pf = request.files.get("portada")
            if pf and pf.filename:
                pext = os.path.splitext(pf.filename)[1].lower()
                if pext in (".png", ".jpg", ".jpeg", ".webp"):
                    try:
                        (BASE_DIR / "static" / "portadas").mkdir(parents=True, exist_ok=True)
                        for oe in (".png", ".jpg", ".jpeg", ".webp"):
                            of = BASE_DIR / "static" / "portadas" / ("pad_" + pack_id + oe)
                            if oe != pext and of.exists():
                                try:
                                    of.unlink()
                                except Exception:
                                    pass
                        pf.save(str(BASE_DIR / "static" / "portadas" / ("pad_" + pack_id + pext)))
                        import time as _t
                        it["portada"] = "portadas/pad_" + pack_id + pext
                        it["portada_ts"] = int(_t.time())
                    except Exception as e:
                        logging.error("portada edit %s: %s", pack_id, e)
            break
        _pads_guardar(lst)
    if regen:
        for i in range(12):
            f = d / ("pad_%d.wav" % i)
            try:
                if f.exists():
                    f.unlink()
            except Exception:
                pass
        _render_pad_async(pack_id)
    flash("Pad actualizado" + (". Regenerando los 12 tonos…" if regen else ""), "success")
    return redirect(url_for("admin_pads"))


@app.route("/admin/pads/<pack_id>/audio/<int:idx>")
@login_required("admin")
def admin_pads_audio(pack_id, idx):
    """Pre-escucha en el admin: sirve el pad renderizado del tono idx (0-11)."""
    if idx < 0 or idx > 11 or not _pad_pack(pack_id):
        abort(404)
    f = CARPETA_PADS / pack_id / ("pad_%d.wav" % idx)
    if not f.is_file():
        abort(404)
    return send_file(str(f))


# ---- API NeuralPlay: pads ----
@app.route("/api/live/perfiles")
def api_live_perfiles():
    """NeuralPlay: roster de perfiles (nombre + acento + prefs) para el visor local. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    return jsonify(usuarios.listar_perfiles(org_actual()))


@app.route("/api/live/pads")
def api_live_pads():
    """NeuralPlay: lista de packs de pads con su estado. Token."""
    if not _token_ok():
        return jsonify({"ok": False, "error": "unauthorized"}), 403
    out = []
    for p in _pads_cargar():
        est = _pad_estado(p["id"])
        out.append({"id": p["id"], "nombre": p.get("nombre"),
                    "base_root": p.get("base_root"), "base_idx": p.get("base_idx", 0),
                    "artista": p.get("artista", ""),
                    "portada": p.get("portada", ""), "portada_ts": p.get("portada_ts", 0),
                    "listo": est["listo"]})
    return jsonify(out)


@app.route("/api/live/pad/<pack_id>/<int:idx>")
def api_live_pad(pack_id, idx):
    """NeuralPlay: sirve el pad (wav) del pack para la raíz idx (0-11). Token."""
    if not _token_ok():
        abort(403)
    if idx < 0 or idx > 11 or not _pad_pack(pack_id):
        abort(404)
    f = CARPETA_PADS / pack_id / ("pad_%d.wav" % idx)
    if not f.is_file():
        abort(404)
    return send_file(str(f))



@app.route("/api/pistas/<int:numero>")
@login_required("musico")
def api_pistas(numero):
    try:
        n = int(request.args.get("t", "0"))
    except ValueError:
        n = 0
    listo = _web_listo(numero, n)
    stems = []
    fam_guardadas = _leer_familias(numero)
    fam_por_stem = {os.path.splitext(k)[0]: v for k, v in fam_guardadas.items()}
    if listo:
        web = _carpeta_tono(numero, n) / "web"
        for f in sorted(web.iterdir()):
            if f.is_file() and f.suffix.lower() == ".mp3":
                stems.append({"name": f.stem, "file": f.name,
                              "familia": fam_por_stem.get(f.stem) or _familia_auto(f.stem)})
    elif _tono_listo(numero, n):
        web = _carpeta_tono(numero, n) / "web"
        if not (web.exists() and (web / ".lock").exists()):
            threading.Thread(target=_asegurar_web, args=(numero, n, org_actual()), daemon=True).start()
    return jsonify({"numero": numero, "tono": n, "listo": listo,
                    "hay_pistas": len(_stems_originales(numero)) > 0,
                    "stems": stems, "secciones": _leer_secciones(numero)})


@app.route("/api/pistas/<int:numero>/render/<n>", methods=["POST"])
@login_required("musico")
def api_render(numero, n):
    try:
        n = int(n)
    except ValueError:
        return jsonify({"error": "tono invalido"}), 400
    if _web_listo(numero, n):
        return jsonify({"listo": True})
    if not _stems_originales(numero):
        return jsonify({"listo": False, "error": "sin pistas"})
    if _tono_listo(numero, n):
        web = _carpeta_tono(numero, n) / "web"
        if not (web.exists() and (web / ".lock").exists()):
            threading.Thread(target=_asegurar_web, args=(numero, n, org_actual()), daemon=True).start()
    else:
        d = _carpeta_tono(numero, n)
        if not (d.exists() and (d / ".lock").exists()):
            threading.Thread(target=_render_tono, args=(numero, n, org_actual()), daemon=True).start()
    return jsonify({"listo": False, "estado": "procesando"})


@app.route("/api/pistas/<int:numero>/render/<n>/estado")
@login_required("musico")
def api_render_estado(numero, n):
    try:
        n = int(n)
    except ValueError:
        return jsonify({"listo": False})
    if _web_listo(numero, n):
        return jsonify({"listo": True})
    prog = ""
    for lock in (_carpeta_tono(numero, n) / ".lock", _carpeta_tono(numero, n) / "web" / ".lock"):
        if lock.exists():
            try:
                prog = lock.read_text().strip()
            except Exception:
                prog = ""
    return jsonify({"listo": False, "progreso": prog})


@app.route("/pista/<int:numero>/<path:archivo>")
@login_required("musico")
def servir_pista(numero, archivo):
    try:
        n = int(request.args.get("t", "0"))
    except ValueError:
        n = 0
    base = (_carpeta_tono(numero, n) / "web").resolve()
    ruta = (base / archivo).resolve()
    try:
        dentro = os.path.commonpath([str(base), str(ruta)]) == str(base)
    except ValueError:
        dentro = False
    if not dentro or not ruta.is_file():
        abort(404)
    return send_file(str(ruta))


# ---- Admin: gestion de pistas de ensayo ----
@app.route("/admin/pistas")
@login_required("admin")
def admin_pistas():
    biblioteca = cargar_biblioteca()
    songs = []
    for numero in sorted(biblioteca.keys()):
        c = biblioteca[numero]
        carpeta = dir_pistas(_cur_org()) / str(numero)
        n = 0
        if carpeta.is_dir():
            n = sum(1 for f in carpeta.iterdir() if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO)
        songs.append({"numero": numero, "titulo": c.get("titulo", ""), "artista": c.get("artista", ""), "n_pistas": n})
    return render_template("admin_pistas.html", songs=songs)


@app.route("/admin/pistas/subir", methods=["POST"])
@login_required("admin")
def admin_pistas_subir():
    numero = request.form.get("numero", "").strip()
    if not numero.isdigit():
        flash("Elegi una cancion valida", "error")
        return redirect(url_for("admin_pistas"))
    archivos = request.files.getlist("pistas")
    carpeta = dir_pistas(_cur_org()) / numero
    carpeta.mkdir(exist_ok=True)
    guardadas = 0
    for a in archivos:
        if not a or a.filename == "":
            continue
        ext = os.path.splitext(a.filename)[1].lower()
        if ext not in _EXT_AUDIO_ENSAYO:
            continue
        nombre = secure_filename(a.filename)
        if not nombre:
            continue
        a.save(str(carpeta / nombre))
        guardadas += 1
    if guardadas:
        _invalidar_tonos(int(numero))
        threading.Thread(target=_asegurar_web, args=(int(numero), 0, org_actual()), daemon=True).start()
        flash("OK: " + str(guardadas) + " pista(s) subida(s) a la cancion #" + numero, "success")
    else:
        flash("No se subio ninguna pista (revisa el formato: mp3/m4a/ogg/wav)", "error")
    return redirect(url_for("admin_pistas"))


@app.route("/admin/pistas/eliminar/<int:numero>", methods=["POST"])
@login_required("admin")
def admin_pistas_eliminar(numero):
    carpeta = dir_pistas(_cur_org()) / str(numero)
    borradas = 0
    if carpeta.is_dir():
        for f in list(carpeta.iterdir()):
            if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO:
                try:
                    f.unlink()
                    borradas += 1
                except Exception:
                    pass
            elif f.is_dir() and (f.name == "web" or f.name.startswith("tono_")):
                try:
                    shutil.rmtree(str(f))
                except Exception:
                    pass
    flash("OK: " + str(borradas) + " pista(s) eliminada(s) de #" + str(numero), "success")
    return redirect(url_for("admin_pistas"))


# ---- Mapa de secciones (etapa 2) ----
def _archivo_secciones(numero):
    return dir_pistas(_cur_org()) / str(numero) / "secciones.json"


def _leer_secciones(numero):
    f = _archivo_secciones(numero)
    if f.is_file():
        try:
            return json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            return []
    return []


def _parse_tiempo(txt):
    txt = txt.strip().replace(",", ".")
    partes = txt.split(":")
    try:
        nums = [float(p) for p in partes]
    except ValueError:
        return None
    if len(nums) == 4:
        h, mi, se, ms = nums; return h*3600 + mi*60 + se + ms/1000.0
    if len(nums) == 3:
        h, mi, se = nums; return h*3600 + mi*60 + se
    if len(nums) == 2:
        mi, se = nums; return mi*60 + se
    if len(nums) == 1:
        return nums[0]
    return None


def _fmt_tiempo(secs):
    ms = int(round((secs - int(secs)) * 1000))
    t = int(secs); h = t // 3600; mi = (t % 3600) // 60; se = t % 60
    return "%02d:%02d:%02d:%03d" % (h, mi, se, ms)


def _hallar_stem_familia(numero, familia):
    base = dir_pistas(_cur_org()) / str(numero)
    if not base.is_dir():
        return None
    fam = _leer_familias(numero)
    for f in sorted(base.iterdir()):
        if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO:
            fa = fam.get(f.name) or _familia_auto(f.stem)
            if fa == familia:
                return f
    return None


def _admin_audio_dir(numero):
    d = dir_pistas(_cur_org()) / str(numero) / "_admin"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _stems_para_mezcla(numero):
    """Todos los stems menos Click y Guia (los instrumentos, para la mezcla de preview)."""
    base = dir_pistas(_cur_org()) / str(numero)
    if not base.is_dir():
        return []
    fam = _leer_familias(numero)
    outs = []
    for f in sorted(base.iterdir()):
        if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO:
            fa = fam.get(f.name) or _familia_auto(f.stem)
            if fa not in ("Click", "Guía"):
                outs.append(f)
    return outs


def _asegurar_audio_admin(numero, tipo, org=None):
    """Genera (si falta) el mp3 de 'guia' o 'mezcla' para el editor de secciones."""
    if org is not None:
        _set_org_ctx(org)
    d = _admin_audio_dir(numero)
    out = d / (tipo + ".mp3")
    lock = d / (tipo + ".lock")
    if out.exists() or lock.exists():
        return
    try:
        lock.write_text("1")
        if tipo == "guia":
            g = _hallar_stem_familia(numero, "Guía")
            if g is None:
                return
            subprocess.run(["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y", "-i", str(g),
                            "-b:a", "128k", str(out)],
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=300)
        elif tipo == "click":
            g = _hallar_stem_familia(numero, "Click")
            if g is None:
                return
            subprocess.run(["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y", "-i", str(g),
                            "-b:a", "128k", str(out)],
                           stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=300)
        elif tipo == "mezcla":
            ins = _stems_para_mezcla(numero)
            if not ins:
                return
            cmd = ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y"]
            for f in ins:
                cmd += ["-i", str(f)]
            cmd += ["-filter_complex",
                    "amix=inputs=%d:normalize=0,alimiter=limit=0.95" % len(ins),
                    "-b:a", "128k", str(out)]
            subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=600)
        elif tipo == "ambas":
            g = _hallar_stem_familia(numero, "Guía")
            ins = ([g] if g else []) + _stems_para_mezcla(numero)
            if not ins:
                return
            cmd = ["/usr/bin/nice", "-n", "19", "/usr/bin/ffmpeg", "-y"]
            for fpath in ins:
                cmd += ["-i", str(fpath)]
            pesos = " ".join((["1.6"] if g else []) + ["1"] * (len(ins) - (1 if g else 0)))
            cmd += ["-filter_complex",
                    "amix=inputs=%d:normalize=0:weights=%s,alimiter=limit=0.95" % (len(ins), pesos),
                    "-b:a", "128k", str(out)]
            subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=600)
    except Exception as e:
        logging.error("audio_admin %s/%s: %s", numero, tipo, e)
    finally:
        try:
            lock.unlink()
        except Exception:
            pass


@app.route("/admin/pistas/<int:numero>/audio_prep/<tipo>", methods=["POST"])
@login_required("admin")
def admin_audio_prep(numero, tipo):
    if tipo not in ("guia", "mezcla", "ambas", "click"):
        return jsonify({"ok": False, "error": "tipo"})
    if tipo == "guia" and _hallar_stem_familia(numero, "Guía") is None:
        return jsonify({"ok": False, "error": "sin_guia"})
    if tipo == "click" and _hallar_stem_familia(numero, "Click") is None:
        return jsonify({"ok": False, "error": "sin_click"})
    if tipo == "mezcla" and not _stems_para_mezcla(numero):
        return jsonify({"ok": False, "error": "sin_stems"})
    if tipo == "ambas" and _hallar_stem_familia(numero, "Guía") is None and not _stems_para_mezcla(numero):
        return jsonify({"ok": False, "error": "sin_stems"})
    out = _admin_audio_dir(numero) / (tipo + ".mp3")
    if out.exists():
        return jsonify({"ok": True, "ready": True})
    lock = _admin_audio_dir(numero) / (tipo + ".lock")
    if not lock.exists():
        threading.Thread(target=_asegurar_audio_admin, args=(numero, tipo, org_actual()), daemon=True).start()
    return jsonify({"ok": True, "ready": False, "generando": True})


@app.route("/admin/pistas/<int:numero>/audio/<tipo>")
@login_required("admin")
def admin_audio(numero, tipo):
    if tipo not in ("guia", "mezcla", "ambas", "click"):
        return ("no", 404)
    out = _admin_audio_dir(numero) / (tipo + ".mp3")
    if not out.exists():
        return ("no listo", 404)
    return send_file(str(out), mimetype="audio/mpeg", conditional=True)


def _hallar_stem_click(numero):
    """Devuelve el archivo del stem de Click (familia 'Click'), o None."""
    base = dir_pistas(_cur_org()) / str(numero)
    if not base.is_dir():
        return None
    fam = _leer_familias(numero)
    cand = []
    for f in sorted(base.iterdir()):
        if f.is_file() and f.suffix.lower() in _EXT_AUDIO_ENSAYO:
            fa = fam.get(f.name) or _familia_auto(f.stem)
            if fa == "Click":
                cand.append(f)
    for f in cand:
        if re.search(r"click|metr", f.stem, re.I):
            return f
    return cand[0] if cand else None


def _detectar_beats_click(numero):
    """Detecta beats y mapa de tempo escuchando el stem de Click.
    Devuelve dict con beats[], tempo[], bpm_prom, wave[] (para dibujar)."""
    try:
        import numpy as np
    except Exception:
        return {"ok": False, "error": "numpy"}
    stem = _hallar_stem_click(numero)
    if stem is None:
        return {"ok": False, "error": "sin_click"}
    sr = 8000
    try:
        p = subprocess.run(
            ["/usr/bin/ffmpeg", "-v", "quiet", "-i", str(stem),
             "-ac", "1", "-ar", str(sr), "-f", "f32le", "-"],
            stdout=subprocess.PIPE, stderr=subprocess.DEVNULL, timeout=120)
        x = np.frombuffer(p.stdout, dtype=np.float32).astype(np.float32)
    except Exception as e:
        logging.error("detectar_click ffmpeg %s: %s", numero, e)
        return {"ok": False, "error": "ffmpeg"}
    if x.size < sr:
        return {"ok": False, "error": "audio_corto"}
    dur = x.size / float(sr)
    # envelope de energia (media movil sobre |x|) + onset = diferencia positiva
    ax = np.abs(x)
    win = max(1, int(sr * 0.005))
    env = np.convolve(ax, np.ones(win, dtype=np.float32) / win, mode="same")
    d = np.diff(env, prepend=env[:1])
    d[d < 0] = 0.0
    mxd = float(d.max()) if d.size else 0.0
    if mxd > 0:
        d = d / mxd
    thr = max(0.14, float(d.mean() + 1.5 * d.std()))
    min_gap = int(sr * 0.11)   # hasta ~545 BPM; evita doble-disparo por un mismo click
    idx = np.where(d > thr)[0]
    beats = []
    fuerzas = []
    if idx.size:
        cortes = np.where(np.diff(idx) > min_gap)[0] + 1
        for g in np.split(idx, cortes):
            jstar = int(g[int(np.argmax(d[g]))])
            beats.append(round(jstar / float(sr), 4))
            fuerzas.append(round(float(d[jstar]), 3))
    # mapa de tempo por segmentos de BPM parecido (+-6%)
    tempo = []
    bpm_prom = 0
    if len(beats) >= 2:
        difs = np.diff(np.array(beats))
        bpms = 60.0 / np.clip(difs, 1e-3, None)
        seg_start = beats[0]
        acc = [float(bpms[0])]
        for i2 in range(1, len(bpms)):
            med = float(np.median(acc))
            if abs(float(bpms[i2]) - med) / max(1.0, med) > 0.06:
                tempo.append({"t": round(float(seg_start), 3), "bpm": int(round(med))})
                seg_start = beats[i2]
                acc = [float(bpms[i2])]
            else:
                acc.append(float(bpms[i2]))
        tempo.append({"t": round(float(seg_start), 3), "bpm": int(round(float(np.median(acc))))})
        bpm_prom = int(round(float(np.median(bpms))))
    # forma de onda comprimida (~700 picos) para dibujar
    W = 700
    step = max(1, x.size // W)
    wave = []
    for i3 in range(0, x.size, step):
        wave.append(float(np.abs(x[i3:i3 + step]).max()))
        if len(wave) >= W:
            break
    mxw = max(wave) if wave else 1.0
    if mxw > 0:
        wave = [round(v / mxw, 3) for v in wave]
    return {"ok": True, "dur": round(dur, 3), "stem": stem.name,
            "beats": beats, "strengths": fuerzas, "tempo": tempo,
            "bpm_prom": bpm_prom, "wave": wave}


@app.route("/admin/pistas/<int:numero>/detectar_click", methods=["POST"])
@login_required("admin")
def admin_detectar_click(numero):
    try:
        res = _detectar_beats_click(numero)
    except Exception as e:
        logging.error("detectar_click %s: %s", numero, e)
        return jsonify({"ok": False, "error": "interno"})
    return jsonify(res)


@app.route("/admin/pistas/<int:numero>/secciones", methods=["GET", "POST"])
@login_required("admin")
def admin_secciones(numero):
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    if not cancion:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    secs_chart = cancion.get("secciones", [])
    if request.method == "POST":
        carpeta = dir_pistas(_cur_org()) / str(numero)
        carpeta.mkdir(exist_ok=True)
        sj = request.form.get("secciones_json")
        if sj is not None:
            # modelo nuevo: el editor manda toda la lista (nombre + t)
            try:
                arr = json.loads(sj)
                if not isinstance(arr, list):
                    arr = []
            except Exception:
                arr = []
            secs = []
            for it in arr:
                if not isinstance(it, dict):
                    continue
                nombre = str(it.get("nombre", "")).strip() or "Seccion"
                t = it.get("t", None)
                try:
                    t = round(float(t), 3)
                except Exception:
                    t = None
                secs.append({"nombre": nombre, "t": t})
            secs.sort(key=lambda s: (s["t"] is None, s["t"] if s["t"] is not None else 0.0))
            guardadas = [{"i": i, "nombre": s["nombre"], "t": s["t"]}
                         for i, s in enumerate(secs) if s["t"] is not None]
            _archivo_secciones(numero).write_text(
                json.dumps(guardadas, ensure_ascii=False, indent=2), encoding="utf-8")
            # sembrar cancion["secciones"] (nombres/orden) preservando acordes por indice
            fdoc = _archivo_cancion(numero)
            if fdoc:
                try:
                    doc = json.loads(fdoc.read_text(encoding="utf-8"))
                except Exception:
                    doc = {}
                viejas = doc.get("secciones", [])
                if not isinstance(viejas, list):
                    viejas = []
                nuevas = []
                for i, s in enumerate(secs):
                    base = {"tipo": s["nombre"]}
                    if i < len(viejas) and isinstance(viejas[i], dict):
                        v = viejas[i]
                        for k in ("nota", "lines", "inst", "prog"):
                            if k in v:
                                base[k] = v[k]
                    nuevas.append(base)
                doc["secciones"] = nuevas
                if not request.form.get("autosave"):
                    _backup_cancion(fdoc)
                fdoc.write_text(json.dumps(doc, ensure_ascii=False, indent=2), encoding="utf-8")
        else:
            # modelo viejo: tiempos t_<i> contra el chart existente
            guardadas = []
            for i, sec in enumerate(secs_chart):
                tiempo = request.form.get("t_" + str(i), "").strip()
                if not tiempo:
                    continue
                val = _parse_tiempo(tiempo)
                if val is None:
                    continue
                guardadas.append({"i": i, "nombre": sec.get("tipo", "Seccion"), "t": round(val, 3)})
            guardadas.sort(key=lambda x: x["t"])
            _archivo_secciones(numero).write_text(
                json.dumps(guardadas, ensure_ascii=False, indent=2), encoding="utf-8")
        if request.form.get("autosave"):
            return jsonify({"ok": True})
        if request.form.get("accion") == "exportar_midi":
            data = _construir_midi_secciones(numero)
            fn = _nombre_base_export(cancion, numero) + "_CHART.mid"
            tmp = Path("/tmp") / fn
            tmp.write_bytes(data)
            return send_file(str(tmp), as_attachment=True, download_name=fn, mimetype="audio/midi")
        if request.form.get("wizard"):
            flash("Cambios guardados", "success")
            return redirect(url_for("admin"))
        flash("Cambios guardados", "success")
        if request.form.get("volver") == "editar":
            return redirect(url_for("admin_editar", numero=numero) + "?tab=tiempos")
        return redirect(url_for("admin_secciones", numero=numero))
    guardadas = {sec["i"]: sec["t"] for sec in _leer_secciones(numero) if "i" in sec}
    filas = []
    for i, sec in enumerate(secs_chart):
        filas.append({"i": i, "tipo": sec.get("tipo", "Seccion"), "nota": sec.get("nota", ""),
                      "t": _fmt_tiempo(guardadas[i]) if i in guardadas else ""})
    titulo = cancion.get("titulo", "Cancion #" + str(numero))
    return render_template("admin_secciones.html", numero=numero, titulo=titulo, filas=filas,
                           wizard=bool(request.args.get("wizard")))


@app.route("/admin/nueva_blanca", methods=["POST"])
@login_required("admin")
def admin_nueva_blanca():
    bib = cargar_biblioteca()
    nuevo = (max(bib.keys()) + 1) if bib else 1
    datos = {"numero": nuevo, "titulo": "Nueva cancion", "artista": "",
             "tono": "", "tempo": 120, "compas": "4/4", "secciones": []}
    destino = dir_canciones(org_actual()) / ("cancion_%03d_nueva.json" % nuevo)
    destino.write_text(json.dumps(datos, ensure_ascii=False, indent=2), encoding="utf-8")
    flash("Cancion nueva creada. Completa su informacion.", "success")
    return redirect(url_for("admin_editar", numero=nuevo) + "?tab=info")


@app.route("/admin/nueva", methods=["GET", "POST"])
@login_required("admin")
def admin_nueva():
    if request.method == "POST":
        archivo = request.files.get("archivo")
        if not archivo or archivo.filename == "":
            flash("Selecciona el archivo JSON", "error")
            return redirect(url_for("admin_nueva"))
        if not archivo.filename.endswith(".json"):
            flash("Solo se aceptan archivos .json", "error")
            return redirect(url_for("admin_nueva"))
        try:
            contenido = archivo.read()
            datos = json.loads(contenido.decode("utf-8"))
            if "numero" not in datos or "titulo" not in datos or "secciones" not in datos:
                flash("El JSON no tiene los campos requeridos (numero, titulo, secciones)", "error")
                return redirect(url_for("admin_nueva"))
        except Exception:
            flash("Archivo no es JSON valido", "error")
            return redirect(url_for("admin_nueva"))
        (dir_canciones(org_actual()) / secure_filename(archivo.filename)).write_bytes(contenido)
        flash("Chart cargado: " + str(datos["titulo"]), "success")
        return redirect(url_for("admin_nueva_pistas", numero=int(datos["numero"])))
    return render_template("admin_nueva.html")


@app.route("/admin/nueva/<int:numero>/pistas", methods=["GET", "POST"])
@login_required("admin")
def admin_nueva_pistas(numero):
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    titulo = cancion.get("titulo", "Cancion #" + str(numero)) if cancion else "Cancion #" + str(numero)
    if request.method == "POST":
        carpeta = dir_pistas(_cur_org()) / str(numero)
        carpeta.mkdir(exist_ok=True)
        guardadas = 0
        for a in request.files.getlist("pistas"):
            if not a or a.filename == "":
                continue
            if os.path.splitext(a.filename)[1].lower() not in _EXT_AUDIO_ENSAYO:
                continue
            nombre = secure_filename(a.filename)
            if not nombre:
                continue
            a.save(str(carpeta / nombre))
            guardadas += 1
        if guardadas:
            _invalidar_tonos(numero)
            threading.Thread(target=_asegurar_web, args=(numero, 0, org_actual()), daemon=True).start()
        flash("OK: " + str(guardadas) + " pista(s) subida(s)", "success")
        return redirect(url_for("admin_secciones", numero=numero, wizard=1))
    return render_template("admin_nueva_pistas.html", numero=numero, titulo=titulo)


FAMILIAS = ["Voces", "AG", "GE", "Piano", "Teclados", "Pad",
            "Cuerdas", "Metales", "Bajo", "Batería", "Percusión", "Loops", "FX",
            "Guía", "Música original", "Click", "Otros"]


def _archivo_midi(numero):
    return dir_pistas(_cur_org()) / str(numero) / "midi.json"


MIDI_CAJAS = [
    ("lyrics", "Lyrics", 16), ("lights1", "Lights 1", 1), ("lights2", "Lights 2", 2),
    ("patches1", "Patches 1", 3), ("patches2", "Patches 2", 4), ("guitar", "Guitar", 5),
    ("aux1", "Aux 1", 6),
]

def _midi_default():
    return {"version": 2, "cajas": {cid: {"canal": ch, "notas": []} for cid, nom, ch in MIDI_CAJAS}}

def _leer_midi(numero):
    f = _archivo_midi(numero)
    data = None
    if f.is_file():
        try:
            data = json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            data = None
    base = _midi_default()
    if isinstance(data, list):                       # formato viejo (una sola lista) -> Lyrics
        base["cajas"]["lyrics"]["notas"] = data
    elif isinstance(data, dict) and isinstance(data.get("cajas"), dict):
        for cid, nom, ch in MIDI_CAJAS:
            c = data["cajas"].get(cid)
            if isinstance(c, dict):
                try:
                    base["cajas"][cid]["canal"] = max(1, min(16, int(c.get("canal", ch))))
                except Exception:
                    pass
                if isinstance(c.get("notas"), list):
                    base["cajas"][cid]["notas"] = c["notas"]
    return base


_NOTA_IDX = {"C": 0, "C#": 1, "Db": 1, "D": 2, "D#": 3, "Eb": 3, "E": 4, "F": 5,
             "F#": 6, "Gb": 6, "G": 7, "G#": 8, "Ab": 8, "A": 9, "A#": 10, "Bb": 10, "B": 11}


def _midi_varlen(n):
    n = int(n)
    if n < 0:
        n = 0
    buf = bytearray([n & 0x7F])
    n >>= 7
    while n:
        buf.append((n & 0x7F) | 0x80)
        n >>= 7
    buf.reverse()
    return bytes(buf)


def _midi_note_num(nota, octava):
    idx = _NOTA_IDX.get(str(nota), 0)
    try:
        oc = int(octava)
    except Exception:
        oc = 0
    return max(0, min(127, (oc + 2) * 12 + idx))  # +1 octava: alineado con ProPresenter/Playback


def _construir_midi(eventos, canal=1):
    TPB = 480
    def s2t(seg):
        return int(round(float(seg) * TPB * 2))  # 120 BPM
    abs_ev = []
    for e in eventos:
        if e.get("desactivar"):
            continue
        seg = e.get("_seg")
        if seg is None:
            continue
        note = _midi_note_num(e.get("nota", "C"), e.get("octava", 0))
        try:
            vel = int(e.get("velocidad", 100))
        except Exception:
            vel = 100
        vel = max(1, min(127, vel))
        on = s2t(seg)
        st = (int(canal) - 1) & 0x0F
        abs_ev.append((on, 0x90 | st, note, vel))
        abs_ev.append((on + 96, 0x80 | st, note, 0))
    abs_ev.sort(key=lambda x: x[0])
    track = bytearray()
    prev = 0
    for tick, status, note, vel in abs_ev:
        track += _midi_varlen(tick - prev)
        track += bytes([status, note, vel])
        prev = tick
    track += _midi_varlen(0) + bytes([0xFF, 0x2F, 0x00])
    header = b"MThd" + (6).to_bytes(4, "big") + (0).to_bytes(2, "big") + (1).to_bytes(2, "big") + TPB.to_bytes(2, "big")
    return header + b"MTrk" + len(track).to_bytes(4, "big") + bytes(track)


def _construir_midi_secciones(numero):
    """Genera la pista MIDI de secciones (idioma del puente / sistema local).
    Canal 16, Program Change = numero-1, nota fantasma de precarga,
    una nota por seccion en escalera desde C1 (36), nota de fin C0 (24)."""
    cancion = cargar_biblioteca().get(numero, {})
    try:
        tempo_bpm = float(cancion.get("tempo") or 120)
    except Exception:
        tempo_bpm = 120
    if tempo_bpm <= 0:
        tempo_bpm = 120
    try:
        num_c, den_c = str(cancion.get("compas", "4/4")).split("/")
        num_c = int(num_c); den_c = int(den_c)
    except Exception:
        num_c, den_c = 4, 4
    TPB = 480
    CANAL = 15          # canal 16 (0-15)
    NOTA_BASE = 36      # C1 = seccion 1
    NOTA_FIN = 24       # C0 = fin

    def s2t(seg):
        return int(round(float(seg) * TPB * tempo_bpm / 60.0))

    secs = [x for x in _leer_secciones(numero) if "i" in x and "t" in x]
    secs.sort(key=lambda x: x["t"])

    eventos = []  # (tick_abs, bytes)
    prog = max(0, min(127, int(numero) - 1))
    eventos.append((0, bytes([0xC0 | CANAL, prog])))
    tg = 2 * TPB  # beat 3 del compas 1 (precarga)
    eventos.append((tg, bytes([0x90 | CANAL, NOTA_BASE, 80])))
    eventos.append((tg + TPB // 4, bytes([0x80 | CANAL, NOTA_BASE, 0])))
    ultimo = tg
    for x in secs:
        tick = s2t(x["t"])
        nota = max(0, min(127, NOTA_BASE + int(x["i"])))
        eventos.append((tick, bytes([0x90 | CANAL, nota, 100])))
        eventos.append((tick + TPB // 4, bytes([0x80 | CANAL, nota, 0])))
        if tick > ultimo:
            ultimo = tick
    if secs:
        tick_fin = ultimo + 1 * num_c * TPB  # 1 compas despues de la ultima
        eventos.append((tick_fin, bytes([0x90 | CANAL, NOTA_FIN, 100])))
        eventos.append((tick_fin + TPB // 4, bytes([0x80 | CANAL, NOTA_FIN, 0])))
    eventos.sort(key=lambda e: e[0])

    track = bytearray()
    nombre = ("Secciones - " + str(cancion.get("titulo", "")))[:120].encode("utf-8", "ignore")
    track += _midi_varlen(0) + bytes([0xFF, 0x03]) + _midi_varlen(len(nombre)) + nombre
    tempo_us = int(round(60000000.0 / tempo_bpm))
    track += _midi_varlen(0) + bytes([0xFF, 0x51, 0x03]) + tempo_us.to_bytes(3, "big")
    dd = max(0, den_c.bit_length() - 1)
    track += _midi_varlen(0) + bytes([0xFF, 0x58, 0x04, num_c & 0xFF, dd & 0xFF, 24, 8])
    prev = 0
    for tick, data in eventos:
        track += _midi_varlen(tick - prev) + data
        prev = tick
    track += _midi_varlen(TPB) + bytes([0xFF, 0x2F, 0x00])
    header = b"MThd" + (6).to_bytes(4, "big") + (0).to_bytes(2, "big") + (1).to_bytes(2, "big") + TPB.to_bytes(2, "big")
    return header + b"MTrk" + len(track).to_bytes(4, "big") + bytes(track)


def _texto_a_lines(texto):
    """Convierte texto con acordes en corchetes -> lineas de tokens [acorde, letra]."""
    lines = []
    for raw in (texto or "").split("\n"):
        tokens = []
        matches = list(re.finditer(r"\[([^\]]*)\]", raw))
        if not matches:
            if raw != "":
                tokens.append(["", raw])
            lines.append(tokens)
            continue
        if matches[0].start() > 0:
            tokens.append(["", raw[:matches[0].start()]])
        for i, m in enumerate(matches):
            end = matches[i + 1].start() if i + 1 < len(matches) else len(raw)
            tokens.append([m.group(1), raw[m.end():end]])
        lines.append(tokens)
    return lines


def _lines_a_texto(lines):
    """Convierte lineas de tokens [acorde, letra] -> texto con corchetes."""
    out = []
    for line in (lines or []):
        seg = ""
        for tok in line:
            ch = tok[0] if len(tok) > 0 else ""
            ly = tok[1] if len(tok) > 1 else ""
            if ch:
                seg += "[" + ch + "]"
            seg += ly
        out.append(seg)
    return "\n".join(out)


def _cifrado_para_editor(cancion):
    out = []
    for sec in (cancion or {}).get("secciones", []):
        if sec.get("inst") and sec.get("prog") is not None:
            texto = " ".join("[" + str(c) + "]" for c in sec.get("prog", []))
        else:
            texto = _lines_a_texto(sec.get("lines", []))
        out.append({"tipo": sec.get("tipo", ""), "nota": sec.get("nota", ""), "texto": texto})
    return out


def _archivo_cancion(numero):
    for f in dir_canciones(org_actual()).glob("*.json"):
        try:
            d = json.loads(f.read_text(encoding="utf-8"))
            if int(d.get("numero", -1)) == int(numero):
                return f
        except Exception:
            continue
    return None


def _familia_auto(nombre):
    n = " " + re.sub(r"[_\-.]+", " ", nombre.lower()) + " "
    if re.search(r"click|metr", n):
        return "Click"
    if re.search(r"guide|gu.a| cue |click.?guide", n):
        return "Guía"
    if re.search(r" voz | voces| vocal| coro| lead| choir| vox | bgv|soprano|alto|tenor", n):
        return "Voces"
    if re.search(r"loop", n):
        return "Loops"
    if re.search(r" fx |efx|sfx|efecto|riser|sweep|impact|whoosh|reverse|uplifter|downlifter", n):
        return "FX"
    if re.search(r" bajo|bass", n):
        return "Bajo"
    if re.search(r" pad ", n):
        return "Pad"
    if re.search(r"piano|rhodes|wurli", n):
        return "Piano"
    if re.search(r"teclado| tecla|keys|synth| sint|organo| k\d", n):
        return "Teclados"
    if re.search(r"ac.stic|acou| ag | ga ", n):
        return "AG"
    if re.search(r"el.ctric|electric| eg | ge |guitarra|guit|gtr| g\d", n):
        return "GE"
    if re.search(r"string|cuerda|viol|cello", n):
        return "Cuerdas"
    if re.search(r"trompeta|trumpet|sax|trombon|brass|metal", n):
        return "Metales"
    if re.search(r"bater|drum|kick|snare| hat| tom|platillo|cymbal|beat", n):
        return "Batería"
    if re.search(r"percu| perc|shaker|conga|tambor|clap|pandero", n):
        return "Percusión"
    if re.search(r"original|mezcla|master|banda| full | todo ", n):
        return "Música original"
    return "Otros"


def _leer_familias(numero):
    f = dir_pistas(_cur_org()) / str(numero) / "familias.json"
    if f.is_file():
        try:
            return json.loads(f.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


@app.route("/admin/pistas/<int:numero>/editar")
@login_required("admin")
def admin_editar(numero):
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    if not cancion:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    fam_guardadas = _leer_familias(numero)
    lista = []
    for f in _stems_originales(numero):
        base = f.rsplit(".", 1)[0]
        lista.append({"file": f, "name": base, "familia": fam_guardadas.get(f) or _familia_auto(base)})
    from transposicion import transponer_acorde, usar_sostenidos
    tono_base = cancion.get("tono", "")
    def _nombre_tono(n):
        if not tono_base:
            return ("+" if n > 0 else "") + str(n)
        tsost = transponer_acorde(tono_base, n, usar_sost=True)
        return transponer_acorde(tono_base, n, usar_sost=usar_sostenidos(tsost))
    tonos = [{"n": n, "listo": _tono_listo(numero, n), "en_proceso": _tono_en_proceso(numero, n), "nombre": _nombre_tono(n)}
             for n in [-5, -4, -3, -2, -1, 1, 2, 3, 4, 5]]
    guardadas_t = {sec["i"]: sec["t"] for sec in _leer_secciones(numero) if "i" in sec}
    filas = []
    for i, sec in enumerate(cancion.get("secciones", [])):
        filas.append({"i": i, "tipo": sec.get("tipo", "Seccion"), "nota": sec.get("nota", ""),
                      "t": _fmt_tiempo(guardadas_t[i]) if i in guardadas_t else ""})
    return render_template("admin_editar.html", numero=numero, titulo=cancion.get("titulo", ""),
                           stems=lista, familias=FAMILIAS, tonos=tonos, filas=filas,
                           midi=_leer_midi(numero),
                           cifrado=_cifrado_para_editor(cancion),
                           letras=_letras_por_seccion(cancion),
                           info={"titulo": cancion.get("titulo", ""), "artista": cancion.get("artista", ""),
                                 "album": cancion.get("album", ""), "genero": cancion.get("genero", ""),
                                 "tema": cancion.get("tema", ""), "tono": cancion.get("tono", ""),
                                 "tempo": cancion.get("tempo", ""), "compas": cancion.get("compas", ""),
                                 "portada": cancion.get("portada", ""), "portada_ts": cancion.get("portada_ts", "")},
                           n_secciones=len(_leer_secciones(numero)))


# ───────── Exportar cifrado a PDF (estilo chart) ─────────
import re as _re_pdf
import html as _htmlmod
from transposicion import PATRON_ACORDE as _PAT_ACORDE

_ABBR_SECCION = [
    (("pre-coro", "pre coro", "precoro", "pre-cor"), "Pr", "#C08A2E"),
    (("post-coro", "post coro", "postcoro", "post-cor"), "Pc", "#C08A2E"),
    (("intro",), "I", "#2E9E8F"),
    (("interludio",), "It", "#7A5AA6"),
    (("instrumental",), "Is", "#2E9E8F"),
    (("verso",), "V", "#3B6FB5"),
    (("coro",), "C", "#B23B4E"),
    (("puente",), "P", "#7A5AA6"),
    (("refra", "refrain"), "Rf", "#2E8B57"),
    (("vamp",), "Vp", "#556070"),
    (("solo",), "So", "#556070"),
    (("tag",), "Tg", "#556070"),
    (("outro",), "O", "#B23B4E"),
    (("final",), "F", "#556070"),
    (("baja",), "Bj", "#556070"),
    (("canales",), "Cn", "#556070"),
    (("exhort",), "Ex", "#556070"),
    (("rap",), "Ra", "#556070"),
    (("acape", "a cape"), "Ac", "#556070"),
    (("pad",), "Pd", "#556070"),
    (("clic", "click"), "Cl", "#556070"),
    (("repet",), "Rp", "#556070"),
]

def _badge_seccion(tipo):
    t = (tipo or "").strip().lower()
    m = _re_pdf.search(r"(\d+)\s*$", t)
    num = m.group(1) if m else ""
    for claves, abbr, color in _ABBR_SECCION:
        if any(t.startswith(k) for k in claves):
            return (abbr + num, color, (tipo or "").strip().upper())
    base = (t[:2] or "S").capitalize()
    return (base + num, "#6b7280", (tipo or "Seccion").strip().upper())

def _acorde_html(ch):
    ch = (ch or "").strip()
    if not ch:
        return ""
    if " " in ch:
        return " ".join(_acorde_html(x) for x in ch.split(" ") if x.strip())
    m = _PAT_ACORDE.match(ch)
    if not m:
        return _htmlmod.escape(ch)
    raiz, modif, bajo = m.groups()
    out = _htmlmod.escape(raiz)
    if modif:
        if modif.startswith("m") and not modif.startswith("maj"):
            out += "m"
            resto = modif[1:]
            if resto:
                out += "<sup>" + _htmlmod.escape(resto) + "</sup>"
        else:
            out += "<sup>" + _htmlmod.escape(modif) + "</sup>"
    if bajo:
        out += "/" + _htmlmod.escape(bajo)
    return out

@app.template_filter("acorde")
def _jinja_acorde(ch):
    from markupsafe import Markup
    return Markup(_acorde_html(ch))

@app.route("/admin/pistas/<int:numero>/pdf")
@login_required("admin")
def admin_cancion_pdf(numero):
    from flask import send_file
    import copy as _copy
    import io as _io
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    if not cancion:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    cancion = _copy.deepcopy(cancion)
    tono_orig = cancion.get("tono", "C")
    tono_pedido = (request.args.get("tono") or "").strip()
    if tono_pedido and tono_pedido != tono_orig:
        from transposicion import NOTA_A_INDICE as _NAI
        mo = _re_pdf.match(r"^([A-G][#b]?)", tono_orig or "")
        md = _re_pdf.match(r"^([A-G][#b]?)", tono_pedido or "")
        ro = mo.group(1) if mo else None
        rd = md.group(1) if md else None
        if ro and rd and ro in _NAI and rd in _NAI:
            sem = (_NAI[rd] - _NAI[ro]) % 12
            if sem:
                cancion = transponer_cancion(cancion, sem)
    secs, mapa = [], []
    for sec in cancion.get("secciones", []):
        abbr, color, nombre = _badge_seccion(sec.get("tipo", ""))
        mapa.append({"abbr": abbr, "color": color})
        lineas = []
        for line in (sec.get("lines") or []):
            toks = [((tk[0] if len(tk) > 0 else ""), (tk[1] if len(tk) > 1 else "")) for tk in line]
            solo = (any((c or "").strip() for c, l in toks) and not any((l or "").strip() for c, l in toks))
            lineas.append({"solo": solo, "toks": toks})
        secs.append({"abbr": abbr, "color": color, "nombre": nombre,
                     "nota": sec.get("nota", "") or "",
                     "inst": bool(sec.get("inst")),
                     "prog": sec.get("prog", []) or [],
                     "lineas": lineas})
    html_doc = render_template("chart_pdf.html", c=cancion, secs=secs, mapa=mapa,
                               tono=cancion.get("tono", tono_orig))
    from weasyprint import HTML as _WEASY_HTML
    pdf_bytes = _WEASY_HTML(string=html_doc).write_pdf()
    titulo = cancion.get("titulo", "cancion")
    artista = cancion.get("artista", "")
    tono_f = cancion.get("tono", tono_orig)
    nombre_arch = titulo + ((" - " + artista) if artista else "") + " (" + str(tono_f) + ").pdf"
    return send_file(_io.BytesIO(pdf_bytes), mimetype="application/pdf",
                     as_attachment=True, download_name=nombre_arch)
# ───────── fin exportar PDF ─────────


@app.route("/admin/pistas/<int:numero>/familias", methods=["POST"])
@login_required("admin")
def admin_familias(numero):
    fam = {}
    for f in _stems_originales(numero):
        v = request.form.get("fam_" + f, "").strip()
        if v:
            fam[f] = v
    (dir_pistas(_cur_org()) / str(numero) / "familias.json").write_text(
        json.dumps(fam, ensure_ascii=False, indent=2), encoding="utf-8")
    flash("Familias guardadas", "success")
    return redirect(url_for("admin_editar", numero=numero))


@app.route("/admin/pistas/<int:numero>/agregar", methods=["POST"])
@login_required("admin")
def admin_editar_subir(numero):
    carpeta = dir_pistas(_cur_org()) / str(numero)
    carpeta.mkdir(exist_ok=True)
    guardadas = 0
    for a in request.files.getlist("pistas"):
        if not a or a.filename == "":
            continue
        if os.path.splitext(a.filename)[1].lower() not in _EXT_AUDIO_ENSAYO:
            continue
        nombre = secure_filename(a.filename)
        if not nombre:
            continue
        a.save(str(carpeta / nombre))
        guardadas += 1
    if guardadas:
        _invalidar_tonos(numero)
        threading.Thread(target=_asegurar_web, args=(numero, 0, org_actual()), daemon=True).start()
    flash("OK: " + str(guardadas) + " pista(s) agregada(s)", "success")
    return redirect(url_for("admin_editar", numero=numero))


@app.route("/admin/pistas/<int:numero>/agregar_uno", methods=["POST"])
@login_required("admin")
def admin_editar_subir_uno(numero):
    a = request.files.get("pista")
    if not a or a.filename == "":
        return jsonify({"ok": False, "error": "sin archivo"}), 400
    if os.path.splitext(a.filename)[1].lower() not in _EXT_AUDIO_ENSAYO:
        return jsonify({"ok": False, "error": "formato"}), 400
    nombre = secure_filename(a.filename)
    if not nombre:
        return jsonify({"ok": False, "error": "nombre"}), 400
    carpeta = dir_pistas(_cur_org()) / str(numero)
    carpeta.mkdir(exist_ok=True)
    a.save(str(carpeta / nombre))
    return jsonify({"ok": True, "name": nombre})


@app.route("/admin/pistas/<int:numero>/subir_fin", methods=["POST"])
@login_required("admin")
def admin_editar_subir_fin(numero):
    _invalidar_tonos(numero)
    threading.Thread(target=_asegurar_web, args=(numero, 0, org_actual()), daemon=True).start()
    return jsonify({"ok": True})


@app.route("/admin/pistas/<int:numero>/borrar_una", methods=["POST"])
@login_required("admin")
def admin_pista_borrar_una(numero):
    archivo = secure_filename(request.form.get("archivo", ""))
    if archivo:
        stem = os.path.splitext(archivo)[0]
        base = dir_pistas(_cur_org()) / str(numero)
        objetivos = [base / archivo, base / "web" / (stem + ".mp3")]
        for d in base.glob("tono_*"):
            if d.is_dir():
                objetivos += [d / archivo, d / "web" / (stem + ".mp3")]
        _invalidar_tonos(numero)
        borro = False
        for p in objetivos:
            if p.is_file():
                try:
                    p.unlink()
                    borro = True
                except Exception:
                    pass
        flash("Pista eliminada: " + archivo if borro else "No se pudo eliminar",
              "success" if borro else "error")
    return redirect(url_for("admin_editar", numero=numero))


# ── Exportar letras (para software de presentacion) ──
def _letras_por_seccion(cancion):
    out = []
    for sec in cancion.get("secciones", []):
        if sec.get("inst"):
            continue
        lineas = []
        for ln in sec.get("lines", []):
            txt = "".join(str(t[1]) for t in ln if isinstance(t, (list, tuple)) and len(t) > 1).rstrip()
            lineas.append(txt)
        while lineas and not lineas[-1].strip():
            lineas.pop()
        if any(x.strip() for x in lineas):
            out.append({"tipo": sec.get("tipo", "Seccion"), "lineas": lineas})
    return out


def _slides_letras(secciones, lineas_por, etiquetas, caso):
    def _c(x):
        if caso == "mayus":
            return x.upper()
        if caso == "minus":
            return x.lower()
        return x
    slides = []
    for sec in secciones:
        lns = [_c(x) for x in sec["lineas"]]
        etq = _c(sec["tipo"]) if etiquetas else None
        if lineas_por and lineas_por > 0:
            for i in range(0, len(lns), lineas_por):
                slides.append({"etiqueta": etq if i == 0 else None, "lineas": lns[i:i + lineas_por]})
        else:
            slides.append({"etiqueta": etq, "lineas": lns})
    return slides


def _letras_txt(slides):
    out = []
    for sl in slides:
        if sl["etiqueta"]:
            out.append(sl["etiqueta"])
        out.extend(sl["lineas"])
        out.append("")
    return (("\n".join(out)).strip() + "\n").encode("utf-8")


def _letras_pptx(slides, fuente, tam, posicion, fondo, colorletra):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from io import BytesIO
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    fg = RGBColor(0, 0, 0) if colorletra == "negro" else RGBColor(0xff, 0xff, 0xff)
    for sl in slides:
        sd = prs.slides.add_slide(blank)
        if fondo == "negro":
            sd.background.fill.solid(); sd.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
        elif fondo == "blanco":
            sd.background.fill.solid(); sd.background.fill.fore_color.rgb = RGBColor(0xff, 0xff, 0xff)
        else:
            sd.background.fill.background()  # transparente (sin relleno)
        if posicion == "tercio":
            top, hgt = Inches(4.4), Inches(2.8)
        else:
            top, hgt = Inches(0.5), Inches(6.5)
        tb = sd.shapes.add_textbox(Inches(0.6), top, Inches(12.1), hgt)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        primero = True
        if sl["etiqueta"]:
            p = tf.paragraphs[0]; primero = False
            p.text = sl["etiqueta"]; p.alignment = PP_ALIGN.CENTER
            r = p.runs[0]; r.font.size = Pt(max(12, int(tam * 0.55))); r.font.name = fuente
            r.font.bold = True; r.font.color.rgb = RGBColor(0xC9, 0xA9, 0x6E)
        for ln in sl["lineas"]:
            p = tf.paragraphs[0] if primero else tf.add_paragraph(); primero = False
            p.text = ln if ln else " "; p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(tam); r.font.name = fuente; r.font.color.rgb = fg
    buf = BytesIO(); prs.save(buf); return buf.getvalue()


def _letras_pdf(slides, fuente, tam, posicion, fondo, colorletra):
    from reportlab.pdfgen import canvas
    from io import BytesIO
    W, H = 960.0, 540.0
    fmap = {"Arial": "Helvetica", "Helvetica": "Helvetica", "Verdana": "Helvetica",
            "Georgia": "Times-Roman", "Times New Roman": "Times-Roman", "Times": "Times-Roman",
            "Courier New": "Courier", "Courier": "Courier",
            "Lora": "Times-Roman", "Montserrat": "Helvetica", "Oswald": "Helvetica"}
    font = fmap.get(fuente, "Helvetica")
    fbold = "Times-Bold" if font == "Times-Roman" else (font + "-Bold")
    fgc = (0, 0, 0) if colorletra == "negro" else (1, 1, 1)
    # PDF no tiene transparencia util: usar fondo con contraste segun el color de letra
    if fondo == "negro":
        bgc = (0, 0, 0)
    elif fondo == "blanco":
        bgc = (1, 1, 1)
    else:
        bgc = (0, 0, 0) if colorletra == "blanco" else (1, 1, 1)
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=(W, H))
    for sl in slides:
        c.setFillColorRGB(*bgc); c.rect(0, 0, W, H, fill=1, stroke=0)
        lines = ([sl["etiqueta"]] if sl["etiqueta"] else []) + sl["lineas"]
        lh = tam * 1.4
        cy = H * 0.24 if posicion == "tercio" else H * 0.5
        y = cy + (len(lines) * lh) / 2 - lh * 0.72
        for i, ln in enumerate(lines):
            es_etq = bool(sl["etiqueta"]) and i == 0
            if es_etq:
                c.setFillColorRGB(0xC9/255.0, 0xA9/255.0, 0x6E/255.0); c.setFont(fbold, max(12, int(tam * 0.55)))
            else:
                c.setFillColorRGB(*fgc); c.setFont(font, tam)
            c.drawCentredString(W / 2, y, ln or "")
            y -= lh
        c.showPage()
    c.save(); return buf.getvalue()


@app.route("/admin/pistas/<int:numero>/letras", methods=["POST"])
@login_required("admin")
def admin_letras(numero):
    from io import BytesIO
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    if not cancion:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    fmt = request.form.get("formato", "pptx")
    fuente = (request.form.get("fuente", "Arial") or "Arial")[:40]
    try:
        tam = int(request.form.get("tamano", 40))
    except Exception:
        tam = 40
    tam = max(10, min(120, tam))
    try:
        lineas_por = int(request.form.get("lineas_por", 4))
    except Exception:
        lineas_por = 4
    lineas_por = max(0, min(20, lineas_por))
    etiquetas = request.form.get("etiquetas") == "1"
    caso = request.form.get("caso", "normal")
    if caso not in ("normal", "mayus", "minus"):
        caso = "normal"
    posicion = "tercio" if request.form.get("posicion", "centro") == "tercio" else "centro"
    fondo = request.form.get("fondo", "transparente")
    if fondo not in ("transparente", "negro", "blanco"):
        fondo = "transparente"
    colorletra = "negro" if request.form.get("colorletra", "blanco") == "negro" else "blanco"
    secs = _letras_por_seccion(cancion)
    slides = _slides_letras(secs, lineas_por, etiquetas, caso)
    base = secure_filename(cancion.get("titulo", "") or ("cancion_%d" % numero)) or ("cancion_%d" % numero)
    if fmt == "txt":
        data = _letras_txt(slides); ext = "txt"; mime = "text/plain; charset=utf-8"
    elif fmt == "pdf":
        data = _letras_pdf(slides, fuente, tam, posicion, fondo, colorletra); ext = "pdf"; mime = "application/pdf"
    else:
        data = _letras_pptx(slides, fuente, tam, posicion, fondo, colorletra); ext = "pptx"
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return send_file(BytesIO(data), as_attachment=True, download_name=base + "_letras." + ext, mimetype=mime)


@app.route("/admin/pistas/<int:numero>/midi", methods=["POST"])
@login_required("admin")
def admin_midi(numero):
    biblioteca = cargar_biblioteca()
    cancion = biblioteca.get(numero)
    if not cancion:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    try:
        payload = json.loads(request.form.get("midi_json", "{}"))
    except Exception:
        payload = {}
    cajas_in = payload.get("cajas", {}) if isinstance(payload, dict) else {}

    def _limpiar_notas(lst):
        out = []
        if not isinstance(lst, list):
            return out
        for e in lst:
            if not isinstance(e, dict):
                continue
            try:
                oc = int(e.get("octava", 0))
            except Exception:
                oc = 0
            try:
                vel = int(e.get("velocidad", 100))
            except Exception:
                vel = 100
            out.append({
                "mensaje": "Nota",
                "t": str(e.get("t", "")).strip(),
                "nota": str(e.get("nota", "C")),
                "octava": oc,
                "velocidad": max(0, min(127, vel)),
                "desactivar": bool(e.get("desactivar")),
                "descripcion": str(e.get("descripcion", "")).strip()[:80],
            })
        return out

    salida = {"version": 2, "cajas": {}}
    for cid, nom, chdef in MIDI_CAJAS:
        c = cajas_in.get(cid) if isinstance(cajas_in.get(cid), dict) else {}
        try:
            canal = int(c.get("canal", chdef))
        except Exception:
            canal = chdef
        canal = max(1, min(16, canal))
        salida["cajas"][cid] = {"canal": canal, "notas": _limpiar_notas(c.get("notas"))}

    carpeta = dir_pistas(_cur_org()) / str(numero)
    carpeta.mkdir(exist_ok=True)
    _archivo_midi(numero).write_text(json.dumps(salida, ensure_ascii=False, indent=2), encoding="utf-8")

    if request.form.get("accion") == "exportar":
        cid = request.form.get("caja", "lyrics")
        lane = salida["cajas"].get(cid) or salida["cajas"].get("lyrics")
        nom_caja = dict((c[0], c[1]) for c in MIDI_CAJAS).get(cid, cid)
        for e in lane["notas"]:
            e["_seg"] = _parse_tiempo(e["t"]) if e["t"] else None
        data = _construir_midi(lane["notas"], canal=lane["canal"])
        base = _nombre_base_export(cancion, numero)
        fn = base + "_" + (secure_filename(nom_caja) or cid) + ".mid"
        tmp = Path("/tmp") / fn
        tmp.write_bytes(data)
        return send_file(str(tmp), as_attachment=True, download_name=fn, mimetype="audio/midi")
    flash("MIDI guardado", "success")
    return redirect(url_for("admin_editar", numero=numero) + "?tab=midi")


@app.route("/admin/pistas/<int:numero>/info", methods=["POST"])
@login_required("admin")
def admin_info(numero):
    f = _archivo_cancion(numero)
    if not f:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    try:
        datos = json.loads(f.read_text(encoding="utf-8"))
    except Exception:
        flash("No se pudo leer la cancion", "error")
        return redirect(url_for("admin_editar", numero=numero) + "?tab=info")
    titulo = request.form.get("titulo", "").strip()
    if titulo:
        datos["titulo"] = titulo[:120]
    datos["artista"] = request.form.get("artista", "").strip()[:80]
    datos["album"] = request.form.get("album", "").strip()[:80]
    datos["genero"] = request.form.get("genero", "").strip()[:40]
    datos["tema"] = request.form.get("tema", "").strip()[:40]
    tono = request.form.get("tono", "").strip()
    if tono:
        datos["tono"] = tono[:8]
    tempo = request.form.get("tempo", "").strip()
    if tempo:
        try:
            _t = float(tempo.replace(",", "."))
            datos["tempo"] = int(_t) if _t == int(_t) else round(_t, 3)
        except Exception:
            pass
    compas = request.form.get("compas", "").strip()
    if compas:
        datos["compas"] = compas[:8]
    pf = request.files.get("portada")
    if pf and pf.filename:
        ext = pf.filename.rsplit(".", 1)[-1].lower() if "." in pf.filename else ""
        if ext in ("png", "jpg", "jpeg", "webp"):
            carpeta = BASE_DIR / "static" / "portadas"
            carpeta.mkdir(parents=True, exist_ok=True)
            for viejo in carpeta.glob(str(numero) + ".*"):
                try:
                    viejo.unlink()
                except Exception:
                    pass
            destino = carpeta / (str(numero) + "." + ext)
            try:
                pf.save(str(destino))
                datos["portada"] = "portadas/" + str(numero) + "." + ext
                import time as _t
                datos["portada_ts"] = int(_t.time())
            except Exception as e:
                flash("No se pudo guardar la portada: %s" % e, "error")
        else:
            flash("Formato de portada no valido (PNG, JPG o WEBP)", "error")
    f.write_text(json.dumps(datos, ensure_ascii=False, indent=2), encoding="utf-8")
    flash("Informacion guardada", "success")
    return redirect(url_for("admin_editar", numero=numero) + "?tab=info")


def _backup_cancion(f, keep=15):
    """Guarda un .bak de la cancion antes de sobrescribirla. Conserva las ultimas 'keep'."""
    try:
        import shutil, time
        bdir = BASE_DIR / "backups_canciones"
        bdir.mkdir(exist_ok=True)
        dest = bdir / (f.stem + "." + time.strftime("%Y%m%d_%H%M%S") + ".bak.json")
        shutil.copy2(str(f), str(dest))
        viejos = sorted(bdir.glob(f.stem + ".*.bak.json"))
        for old in viejos[:-keep]:
            try:
                old.unlink()
            except Exception:
                pass
    except Exception as e:
        print("backup cancion fallo:", e)


@app.route("/admin/pistas/<int:numero>/cifrado", methods=["POST"])
@login_required("admin")
def admin_cifrado(numero):
    f = _archivo_cancion(numero)
    if not f:
        flash("Cancion no encontrada", "error")
        return redirect(url_for("admin_pistas"))
    try:
        datos = json.loads(f.read_text(encoding="utf-8"))
    except Exception:
        flash("No se pudo leer la cancion", "error")
        return redirect(url_for("admin_editar", numero=numero) + "?tab=cifrado")
    try:
        secs_in = json.loads(request.form.get("cifrado_json", "[]"))
        if not isinstance(secs_in, list):
            secs_in = []
    except Exception:
        secs_in = []
    nuevas = []
    for it in secs_in:
        if not isinstance(it, dict):
            continue
        tipo = str(it.get("tipo", "")).strip() or "Seccion"
        nota = str(it.get("nota", "")).strip()
        lineas = _texto_a_lines(it.get("texto", ""))
        tiene_letra = any((len(t) > 1 and str(t[1]).strip()) for ln in lineas for t in ln)
        sec = {"tipo": tipo}
        if nota:
            sec["nota"] = nota
        if not tiene_letra and len(lineas) <= 1:
            sec["inst"] = True
            sec["prog"] = [str(t[0]) for ln in lineas for t in ln if len(t) > 0 and str(t[0]).strip()]
        else:
            sec["lines"] = lineas
        nuevas.append(sec)
    datos["secciones"] = nuevas
    if request.form.get("autosave"):
        f.write_text(json.dumps(datos, ensure_ascii=False, indent=2), encoding="utf-8")
        return jsonify({"ok": True})
    _backup_cancion(f)
    f.write_text(json.dumps(datos, ensure_ascii=False, indent=2), encoding="utf-8")
    flash("Cifrado guardado", "success")
    return redirect(url_for("admin_editar", numero=numero) + "?tab=cifrado")


@app.route("/admin/pistas/<int:numero>/borrar_tono/<n>", methods=["POST"])
@login_required("admin")
def admin_borrar_tono(numero, n):
    try:
        n = int(n)
    except ValueError:
        return redirect(url_for("admin_editar", numero=numero) + "?tab=transposicion")
    if n != 0:
        d = _carpeta_tono(numero, n)
        if d.is_dir():
            try:
                shutil.rmtree(str(d))
                flash("Transposición eliminada", "success")
            except Exception as e:
                flash("No se pudo eliminar: %s" % e, "error")
    return redirect(url_for("admin_editar", numero=numero) + "?tab=transposicion")


@app.route("/admin/pistas/<int:numero>/pregenerar", methods=["POST"])
@login_required("admin")
def admin_pregenerar(numero):
    try:
        n = int(request.form.get("tono", "0"))
    except ValueError:
        n = 0
    if n == 0 or not _stems_originales(numero):
        flash("Tono invalido o sin pistas", "error")
        return redirect(url_for("admin_editar", numero=numero))
    if _tono_listo(numero, n):
        flash("El tono " + str(n) + " ya estaba listo", "success")
    else:
        d = _carpeta_tono(numero, n)
        if not (d.exists() and (d / ".lock").exists()):
            threading.Thread(target=_render_tono, args=(numero, n, org_actual()), daemon=True).start()
        flash("Generando el tono " + ("+" if n > 0 else "") + str(n) +
              " en segundo plano. Puede tardar unos minutos; refresca para ver cuando este listo.", "success")
    return redirect(url_for("admin_editar", numero=numero))


if __name__ == "__main__":
    # En desarrollo local (no en VPS), correr en puerto 5051
    # En VPS, esto lo manejará gunicorn o systemd
    app.run(host="127.0.0.1", port=5051, debug=False)
